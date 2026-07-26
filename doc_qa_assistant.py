#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
基于本地 Qwen / 千问 API 的文档问答脚本。
支持常见办公格式：PDF / Word(docx) / Excel(xlsx) / PowerPoint(pptx) / CSV / TXT / MD；
以及图片 PNG / JPG / WebP / GIF 等（RapidOCR 识别图中文字后参与检索）。

流程：解析文件 → 分块并记录页码或等价位置 → 向量检索 → 用大模型根据检索片段作答并引用位置。

环境变量：
  MS_MODEL_ID   模型 ID（魔搭与 HuggingFace 上 Qwen 常用同名，如 Qwen/Qwen2.5-3B-Instruct）
  LLM_HUB       加载渠道：auto（先魔搭，失败则 HF）| modelscope | huggingface
  HF_ENDPOINT   可选；未设置且未 export RAG_USE_OFFICIAL_HF=1 时，默认 https://hf-mirror.com（减轻直连 huggingface.co 的 SSL EOF）
  RAG_USE_OFFICIAL_HF  设为 1/true 时不再默认 HF_ENDPOINT，改走官方 Hub（需网络可达）
  MS_EMBED_ID   可选，嵌入模型；默认 BAAI/bge-small-zh-v1.5
  DASHSCOPE_API_KEY  可选，若提供则直接走千问兼容 API，不加载本地大模型
  QWEN_API_MODEL     可选，API 模型名，默认 qwen-plus
  QWEN_API_BASE      可选，默认 https://dashscope.aliyuncs.com/compatible-mode/v1
  RAG_STRICT_IMAGE_OCR  设为 1/true 时，若未安装 RapidOCR 则上传图片直接报错；默认关闭（无 OCR 时仍入库占位文本，避免崩溃）

用法示例：
  python doc_qa_assistant.py --files ./a.pdf ./b.docx --question "合同金额是多少？"
  python doc_qa_assistant.py --files ./data.xlsx --question "汇总哪一列？" --top-k 5
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import sys
import tempfile
import threading
from dataclasses import dataclass
from pathlib import Path
from typing import Any, List, Optional, Sequence  # noqa: F401 used in _bundle_key payload

# 在首次访问 Hub 之前生效：默认经镜像拉嵌入权重，避免直连 huggingface.co 频繁 SSL EOF 及 httpx 重试异常。
if os.environ.get("RAG_USE_OFFICIAL_HF", "").strip().lower() not in {"1", "true", "yes", "on"}:
    os.environ.setdefault("HF_ENDPOINT", "https://hf-mirror.com")

import numpy as np

# SentenceTransformer 首次会从 Hub 拉权重；多线程并发加载会触发 huggingface_hub 内 httpx 客户端关闭类错误，故进程内单例串行加载。
_st_model_singleton: dict[str, object] = {}
_st_model_singleton_lock = threading.Lock()


def _faiss_write_index(index, dest: Path) -> None:
    """
    将 FAISS 索引写入目标路径。
    Windows 下原生 faiss 对含非 ASCII 的路径可能 fopen 失败，先写入系统临时文件再移动。
    """
    import faiss

    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    fd, tmp = tempfile.mkstemp(suffix=".faissindex")
    os.close(fd)
    tmp_path = Path(tmp)
    try:
        faiss.write_index(index, str(tmp_path))
        shutil.move(str(tmp_path), dest)
    except Exception:
        tmp_path.unlink(missing_ok=True)
        raise


def _faiss_read_index(src: Path):
    import faiss

    src = Path(src)
    fd, tmp = tempfile.mkstemp(suffix=".faissindex")
    os.close(fd)
    tmp_path = Path(tmp)
    try:
        shutil.copy2(src, tmp_path)
        return faiss.read_index(str(tmp_path))
    finally:
        tmp_path.unlink(missing_ok=True)


# ------------- 数据结构与分块 -------------


@dataclass
class TextChunk:
    text: str
    source: str  # 文件名
    page_label: str  # 人类可读的页码或位置说明
    meta: str  # 额外说明（工作表名等）
    doc_path: str = ""  # 原始文件绝对路径
    chunk_id: str = ""  # 用于缓存/去重的稳定 ID
    kb_note_id: str = ""  # 知识库笔记 id（非笔记来源为空）
    kb_attachment_id: str = ""  # 知识库附件 id
    session_file_id: str = ""  # 会话上传文件 id（chromadb）
    context_header: str = ""  # Contextual Retrieval：块所在文档/位置的语境前缀（仅用于检索，不进 prompt）
    parent_id: str = ""  # 父级单元（页/幻灯片/工作表）标识，同页子块共享
    parent_text: str = ""  # 父级单元全文；命中子块后展开给 LLM 作为证据，子块=父级时为空


def _clean_text(s: str) -> str:
    s = s.replace("\x00", " ")
    s = re.sub(r"\s+", " ", s).strip()
    return s


def normalize_block_text(s: str) -> str:
    """规范空白但保留段落与行边界，供分块时寻找自然断点。

    `_clean_text` 会把换行折叠成空格，那样 `_find_split` 的段落/行分支永不命中，
    分块只能退化到句子或空格边界。
    """
    if not s:
        return ""
    s = s.replace("\x00", " ").replace("\r\n", "\n").replace("\r", "\n")
    s = re.sub(r"[ \t\u00a0\u3000\f\v]+", " ", s)
    s = re.sub(r" *\n *", "\n", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


# 结构/版面感知分块：优先在自然边界（段落 > 句子 > 词）切开，而非在 max_chars 处硬切，
# 避免把同一句话/概念拦腰截断导致嵌入与 BM25 语义受损。仅改变"在哪里切"，不改接口与页码。
_SENTENCE_END = re.compile(r"[。．！？；]|[.!?;](?=\s|$)")


def _find_split(text: str, start: int, hard_end: int) -> int:
    """在 [start, hard_end] 内寻找靠近 hard_end 的自然断点；找不到则返回 hard_end。

    回看窗口约为本块长度的 35%，只在该窗口内回退，保证块长不过度缩水、且必然向前推进。
    """
    n = len(text)
    if hard_end >= n:
        return n
    window = max(1, int((hard_end - start) * 0.35))
    lo = max(start + 1, hard_end - window)
    # 1) 段落边界
    p = text.rfind("\n\n", lo, hard_end)
    if p != -1:
        return p + 2
    p = text.rfind("\n", lo, hard_end)
    if p != -1:
        return p + 1
    # 2) 句子边界（中英文标点）
    best = -1
    for m in _SENTENCE_END.finditer(text, lo, hard_end):
        best = m.end()
    if best != -1 and best > start:
        return best
    # 3) 词/空白边界
    p = text.rfind(" ", lo, hard_end)
    if p != -1 and p + 1 > start:
        return p + 1
    return hard_end


def _env_int(name: str, default: int, minimum: int = 1) -> int:
    raw = os.environ.get(name, "").strip()
    if not raw:
        return default
    try:
        return max(minimum, int(float(raw)))
    except ValueError:
        return default


def _env_float(name: str, default: float, minimum: float, maximum: float) -> float:
    raw = os.environ.get(name, "").strip()
    if not raw:
        return default
    try:
        return max(minimum, min(maximum, float(raw)))
    except ValueError:
        return default


# 分块以 token 为口径：嵌入模型、重排器与提示词预算都按 token 计费，按字符计长会让
# 英文课件的块只有中文块的四分之一信息量（900 字符英文约 220 token，中文约 800 token）。
CHUNK_TOKENS = _env_int("RAG_CHUNK_TOKENS", 480, minimum=64)
CHUNK_OVERLAP_RATIO = _env_float("RAG_CHUNK_OVERLAP_RATIO", 0.15, 0.0, 0.5)
# 低于该长度的块信息量太低（页眉、章节标题、单行幻灯片标题），并入相邻块。
MIN_CHUNK_CHARS = _env_int("RAG_MIN_CHUNK_CHARS", 120, minimum=0)
# 父级单元（整页/整张幻灯片）超过该长度就不再随子块缓存，避免提示词被单页表格挤爆。
PARENT_MAX_CHARS = _env_int("RAG_PARENT_MAX_CHARS", 2400, minimum=0)

# 单一配置源：所有解析器都从这里取分块参数，勿在调用点写死数字。
# max_tokens 缺省沿用 CHUNK_TOKENS；表格类内容行与行相对独立，块可以大一些。
CHUNK_CONFIG: dict[str, dict[str, float]] = {
    "_default": {"max_tokens": CHUNK_TOKENS, "overlap_ratio": CHUNK_OVERLAP_RATIO},
    ".xlsx": {"max_tokens": _env_int("RAG_CHUNK_TOKENS_TABULAR", 640, minimum=64)},
    ".csv": {"max_tokens": _env_int("RAG_CHUNK_TOKENS_TABULAR", 640, minimum=64)},
}
CACHE_VERSION = "rag_cache_v6"  # v6: PDF 图片页 OCR 兜底（v5: token 口径分块 + 短块合并去重 + 父子块）

# 讲义常把整页导出为图片（扫描件、手写公式页、截图排版），文字层为空。
# 这类页在检索里等于不存在，因此文字过少时渲染成位图交给 OCR。
PDF_OCR_MIN_CHARS = _env_int("RAG_PDF_OCR_MIN_CHARS", 50, minimum=0)
PDF_OCR_DPI = _env_int("RAG_PDF_OCR_DPI", 180, minimum=72)
# OCR 每页约数秒，且只在首次建索引时发生（结果随向量缓存复用）；上限防止超大扫描件把上传卡死。
PDF_OCR_MAX_PAGES = _env_int("RAG_PDF_OCR_MAX_PAGES", 400, minimum=0)

# token 估算：CJK 基本一字一 token，拉丁文本约 4 字符一 token（BERT/XLM-R 量级）。
# 估算而非真调 tokenizer，是为了让分块结果与"嵌入模型是否已加载/能否联网"无关，
# 保证 chunk_id 与向量缓存 key 稳定可复现。需要精确计数时设 RAG_CHUNK_TOKENIZER。
_CJK_RE = re.compile(r"[\u2e80-\u9fff\uf900-\ufaff\uff00-\uffef\uac00-\ud7af]")
_LATIN_CHARS_PER_TOKEN = 4.0
_tokenizer_cache: dict[str, Any] = {}


def _chunk_tokenizer():
    """可选的精确 tokenizer（RAG_CHUNK_TOKENIZER=模型名或本地目录）；不可用时返回 None。"""
    ref = os.environ.get("RAG_CHUNK_TOKENIZER", "").strip()
    if not ref:
        return None
    if ref in _tokenizer_cache:
        return _tokenizer_cache[ref]
    try:
        from transformers import AutoTokenizer

        _tokenizer_cache[ref] = AutoTokenizer.from_pretrained(ref)
    except Exception as e:  # 缺依赖、权重拉不下来时静默回退到估算
        print(f"[分块] tokenizer {ref} 不可用，改用估算计数: {e!r}", file=sys.stderr)
        _tokenizer_cache[ref] = None
    return _tokenizer_cache[ref]


def count_tokens(text: str) -> float:
    if not text:
        return 0.0
    tok = _chunk_tokenizer()
    if tok is not None:
        try:
            return float(len(tok.encode(text, add_special_tokens=False)))
        except Exception:
            pass
    cjk = len(_CJK_RE.findall(text))
    return cjk + max(0, len(text) - cjk) / _LATIN_CHARS_PER_TOKEN


def token_budget_to_chars(text: str, max_tokens: int) -> int:
    """把 token 预算折算成本段文本的字符预算（按整段实际 token 密度换算）。"""
    if not text:
        return max_tokens
    density = count_tokens(text) / len(text)
    if density <= 0:
        density = 1.0 / _LATIN_CHARS_PER_TOKEN
    return max(64, int(max_tokens / density))


def chunk_params(ext: str) -> tuple[int, float]:
    """返回该扩展名的 (max_tokens, overlap_ratio)。"""
    default = CHUNK_CONFIG["_default"]
    conf = CHUNK_CONFIG.get((ext or "").lower(), {})
    max_tokens = int(conf.get("max_tokens", default["max_tokens"]))
    overlap_ratio = float(conf.get("overlap_ratio", default["overlap_ratio"]))
    return max_tokens, overlap_ratio


def chunk_unit_text(
    text: str,
    source: str,
    page_label: str,
    meta: str,
    ext: str,
) -> List[TextChunk]:
    """把一个"父级单元"（一页 / 一张幻灯片 / 一个工作表）切成子块。

    子块记录所属父级全文，检索命中子块后可在提示词里展开父级（parent-child retrieval）。
    """
    unit = normalize_block_text(text)
    if not unit:
        return []
    max_tokens, overlap_ratio = chunk_params(ext)
    max_chars = token_budget_to_chars(unit, max_tokens)
    overlap = int(max_chars * overlap_ratio)
    pieces = chunk_by_chars(
        unit,
        source=source,
        page_label=page_label,
        meta=meta,
        max_chars=max_chars,
        overlap=overlap,
    )
    for piece in pieces:
        piece.parent_text = unit
    return pieces


def chunk_by_chars(
    text: str,
    source: str,
    page_label: str,
    meta: str,
    max_chars: int,
    overlap: int,
) -> List[TextChunk]:
    text = normalize_block_text(text)
    if not text:
        return []
    overlap = max(0, min(overlap, max_chars - 1))
    chunks: List[TextChunk] = []
    start = 0
    n = len(text)
    while start < n:
        hard_end = min(start + max_chars, n)
        end = _find_split(text, start, hard_end)
        # 兜底：断点未向前推进时退回硬切，杜绝死循环。
        if end <= start:
            end = hard_end
        piece = text[start:end].strip()
        if piece:
            chunks.append(TextChunk(text=piece, source=source, page_label=page_label, meta=meta))
        if end >= n:
            break
        # 重叠回退后仍必须向前推进，否则重叠大于实际块长时会死循环。
        start = max(end - overlap, start + 1)
    return chunks

# Contextual Retrieval（Anthropic, 2024）：给每个块拼一段"它在文档中的语境"再嵌入+建 BM25。
# 这里用零成本的确定性上下文头（文件名/位置/说明），显著改善术语错配的召回；
# 关闭时置空即可（需清缓存或已由 CACHE_VERSION 变更触发重建）。
CONTEXTUAL_HEADERS = os.environ.get("RAG_CONTEXTUAL_HEADERS", "1").strip().lower() not in {"0", "false", "no", "off"}


def _build_context_header(source: str, page_label: str, meta: str) -> str:
    if not CONTEXTUAL_HEADERS:
        return ""
    parts = []
    if source:
        parts.append(f"Document: {source}")
    if page_label:
        parts.append(f"Location: {page_label}")
    if meta:
        parts.append(str(meta)[:120])
    return " | ".join(parts)


def _sha1_text(text: str) -> str:
    return hashlib.sha1(text.encode("utf-8")).hexdigest()


def _normalized_path(path: Path | str) -> str:
    return str(Path(path).expanduser().resolve())


def _file_fingerprint(path: Path) -> str:
    stat = path.stat()
    payload = {
        "path": _normalized_path(path),
        "size": stat.st_size,
        "mtime_ns": stat.st_mtime_ns,
    }
    return _sha1_text(json.dumps(payload, ensure_ascii=False, sort_keys=True))


def _chunk_identity(chunk: TextChunk) -> str:
    base = {
        "doc_path": chunk.doc_path,
        "source": chunk.source,
        "page_label": chunk.page_label,
        "meta": chunk.meta,
        "text": chunk.text,
    }
    return _sha1_text(json.dumps(base, ensure_ascii=False, sort_keys=True))


_PAGE_LABEL_RE = re.compile(r"^第(\d+)(页|张幻灯片)$")


def _merge_page_labels(labels: Sequence[str]) -> str:
    """合并相邻单元后给出如实的位置说明，如「第16-17张幻灯片」。"""
    uniq: List[str] = []
    for label in labels:
        if label and label not in uniq:
            uniq.append(label)
    if len(uniq) <= 1:
        return uniq[0] if uniq else ""
    parsed = [_PAGE_LABEL_RE.match(label) for label in uniq]
    units = {m.group(2) for m in parsed if m}
    if all(parsed) and len(units) == 1:
        numbers = sorted(int(m.group(1)) for m in parsed if m)
        return f"第{numbers[0]}-{numbers[-1]}{units.pop()}"
    return "、".join(uniq[:3])


def _merge_short_chunks(chunks: Sequence[TextChunk], ext: str) -> List[TextChunk]:
    """把过短的块并入相邻块：先在同一页内合并，再把整页级短块并入相邻页。

    页眉、章节标题页、只有一行标题的幻灯片单独成块时，既占检索命中位又占引用位，
    却几乎不含可作答的信息。跨页合并受 token 预算约束，并如实标成「第15-16张幻灯片」。
    """
    if MIN_CHUNK_CHARS <= 0:
        return list(chunks)
    max_tokens, _overlap_ratio = chunk_params(ext)

    merged: List[TextChunk] = []
    for chunk in chunks:
        prev = merged[-1] if merged else None
        same_unit = (
            prev is not None
            and prev.page_label == chunk.page_label
            and prev.meta == chunk.meta
            and prev.source == chunk.source
        )
        short = len(chunk.text) < MIN_CHUNK_CHARS
        prev_short = prev is not None and len(prev.text) < MIN_CHUNK_CHARS
        if same_unit and (short or prev_short):
            prev.text = f"{prev.text}\n{chunk.text}".strip()
            continue
        merged.append(
            TextChunk(
                text=chunk.text,
                source=chunk.source,
                page_label=chunk.page_label,
                meta=chunk.meta,
                parent_text=chunk.parent_text,
            )
        )

    out: List[TextChunk] = []
    labels: List[List[str]] = []
    for chunk in merged:
        prev = out[-1] if out else None
        if prev is not None and prev.meta == chunk.meta and prev.source == chunk.source:
            either_short = (
                len(chunk.text) < MIN_CHUNK_CHARS or len(prev.text) < MIN_CHUNK_CHARS
            )
            combined = f"{prev.text}\n{chunk.text}".strip()
            if either_short and count_tokens(combined) <= max_tokens:
                prev.text = combined
                labels[-1].append(chunk.page_label)
                prev.page_label = _merge_page_labels(labels[-1])
                # 跨页合并后原父级页已被拼进正文，无需再单独展开。
                prev.parent_text = ""
                continue
        out.append(chunk)
        labels.append([chunk.page_label])
    return out


def _dedupe_chunks(chunks: Sequence[TextChunk]) -> List[TextChunk]:
    """同一文档内正文完全相同的块只保留首次出现的位置。"""
    seen: set[str] = set()
    out: List[TextChunk] = []
    for chunk in chunks:
        key = re.sub(r"\s+", " ", chunk.text).strip().lower()
        if not key or key in seen:
            continue
        seen.add(key)
        out.append(chunk)
    return out


def _finalize_chunks(chunks: Sequence[TextChunk], path: Path) -> List[TextChunk]:
    out: List[TextChunk] = []
    doc_path = _normalized_path(path)
    for chunk in chunks:
        item = TextChunk(
            text=chunk.text,
            source=chunk.source,
            page_label=chunk.page_label,
            meta=chunk.meta,
            doc_path=doc_path,
        )
        parent_text = (chunk.parent_text or "").strip()
        # 父级与子块等长说明该页只有一个块，无需重复缓存；过长的父级不进提示词预算。
        if parent_text and parent_text != item.text and len(parent_text) <= PARENT_MAX_CHARS:
            item.parent_text = parent_text
            item.parent_id = _sha1_text(f"{doc_path}\n{item.page_label}\n{parent_text}")
        else:
            item.parent_id = _sha1_text(f"{doc_path}\n{item.page_label}")
        item.context_header = _build_context_header(chunk.source, chunk.page_label, chunk.meta)
        item.chunk_id = _chunk_identity(item)
        out.append(item)
    return out


def _chunk_to_dict(chunk: TextChunk) -> dict:
    return {
        "text": chunk.text,
        "source": chunk.source,
        "page_label": chunk.page_label,
        "meta": chunk.meta,
        "doc_path": chunk.doc_path,
        "chunk_id": chunk.chunk_id,
        "kb_note_id": getattr(chunk, "kb_note_id", "") or "",
        "kb_attachment_id": getattr(chunk, "kb_attachment_id", "") or "",
        "session_file_id": getattr(chunk, "session_file_id", "") or "",
        "context_header": getattr(chunk, "context_header", "") or "",
        "parent_id": getattr(chunk, "parent_id", "") or "",
        "parent_text": getattr(chunk, "parent_text", "") or "",
    }


def _chunk_from_dict(data: dict) -> TextChunk:
    return TextChunk(
        text=data["text"],
        source=data["source"],
        page_label=data["page_label"],
        meta=data["meta"],
        doc_path=data.get("doc_path", ""),
        chunk_id=data.get("chunk_id", ""),
        kb_note_id=data.get("kb_note_id") or "",
        kb_attachment_id=data.get("kb_attachment_id") or "",
        session_file_id=data.get("session_file_id") or "",
        context_header=data.get("context_header") or "",
        parent_id=data.get("parent_id") or "",
        parent_text=data.get("parent_text") or "",
    )


# ------------- 各格式解析 -------------


def _pdf_ocr_enabled() -> bool:
    return os.environ.get("RAG_PDF_OCR", "1").strip().lower() not in {"0", "false", "no"}


def _ocr_pdf_page(page: Any) -> str:
    """把 PDF 页渲染成位图后 OCR；任何失败都返回空串，不影响其余页入库。"""
    tmp_path: Optional[str] = None
    try:
        ocr = _get_rapid_ocr()
        pix = page.get_pixmap(dpi=PDF_OCR_DPI)
        fd, tmp_path = tempfile.mkstemp(suffix=".png")
        os.close(fd)
        Path(tmp_path).write_bytes(pix.tobytes("png"))
        lines = _text_lines_from_rapidocr_output(ocr(tmp_path))
    except Exception:
        return ""
    finally:
        if tmp_path:
            Path(tmp_path).unlink(missing_ok=True)
    return "\n".join(lines).strip()


def parse_pdf(path: Path) -> List[TextChunk]:
    import fitz  # pymupdf

    out: List[TextChunk] = []
    doc = fitz.open(path)
    name = path.name
    ocr_budget = PDF_OCR_MAX_PAGES if _pdf_ocr_enabled() else 0
    try:
        for i in range(len(doc)):
            page = doc.load_page(i)
            text = page.get_text("text") or ""
            meta = "PDF"
            if len(text.strip()) < PDF_OCR_MIN_CHARS and ocr_budget > 0:
                ocr_budget -= 1
                recognized = _ocr_pdf_page(page)
                if len(recognized) > len(text.strip()):
                    text = recognized
                    meta = "PDF 图片页 OCR"
            page_no = i + 1
            label = f"第{page_no}页"
            out.extend(
                chunk_unit_text(
                    text,
                    source=name,
                    page_label=label,
                    meta=meta,
                    ext=".pdf",
                )
            )
    finally:
        doc.close()
    return out


def parse_docx(path: Path) -> List[TextChunk]:
    from docx import Document

    doc = Document(path)
    name = path.name
    paras: List[str] = []
    for p in doc.paragraphs:
        t = normalize_block_text(p.text)
        if t:
            paras.append(t)
    full = "\n\n".join(paras)
    # Word 无稳定物理页码，按字数估算「约第 N 页」（按约 1800 字/页）
    chars_per_page = 1800
    out: List[TextChunk] = []
    pos = 0
    for para in paras:
        est_page = pos // chars_per_page + 1
        label = f"约第{est_page}页(估算)"
        out.extend(
            chunk_unit_text(
                para,
                source=name,
                page_label=label,
                meta="Word(按字数估算页码)",
                ext=".docx",
            )
        )
        pos += len(para) + 1
    if not out and full:
        out.extend(
            chunk_unit_text(
                full,
                source=name,
                page_label="全文",
                meta="Word",
                ext=".docx",
            )
        )
    return out


def parse_xlsx(path: Path) -> List[TextChunk]:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    name = path.name
    out: List[TextChunk] = []
    try:
        for sheet in wb.worksheets:
            rows: List[str] = []
            for r_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if not cells:
                    continue
                line = " | ".join(cells)
                rows.append(f"行{r_idx}: {line}")
            sheet_text = "\n".join(rows)
            label = f"工作表「{sheet.title}」"
            out.extend(
                chunk_unit_text(
                    sheet_text,
                    source=name,
                    page_label=label,
                    meta="Excel(位置为工作表+行)",
                    ext=".xlsx",
                )
            )
    finally:
        wb.close()
    return out


def parse_pptx(path: Path) -> List[TextChunk]:
    from pptx import Presentation

    prs = Presentation(path)
    name = path.name
    out: List[TextChunk] = []
    for si, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(normalize_block_text(shape.text))
        text = "\n".join(part for part in parts if part)
        label = f"第{si}张幻灯片"
        out.extend(
            chunk_unit_text(
                text,
                source=name,
                page_label=label,
                meta="PPT",
                ext=".pptx",
            )
        )
    return out


def parse_csv(path: Path) -> List[TextChunk]:
    import pandas as pd

    name = path.name
    try:
        df = pd.read_csv(path, encoding="utf-8")
    except UnicodeDecodeError:
        df = pd.read_csv(path, encoding="gbk")
    lines = []
    for i, row in df.iterrows():
        lines.append(f"行{int(i)+2}: " + " | ".join(str(x) for x in row.values))  # +2: 表头占1，pandas 0-based
    text = "\n".join(lines)
    return chunk_unit_text(
        text,
        source=name,
        page_label="CSV 行号见各行前缀",
        meta="CSV",
        ext=".csv",
    )


def parse_txt(path: Path) -> List[TextChunk]:
    raw = path.read_bytes()
    for enc in ("utf-8", "gbk", "utf-16"):
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = raw.decode("utf-8", errors="replace")
    name = path.name
    lines = text.splitlines()
    # 按行范围分块并估算页（约40行/页）
    lines_per_page = 40
    out: List[TextChunk] = []
    buf: List[str] = []
    base_line = 0
    for i, line in enumerate(lines):
        buf.append(line)
        if len(buf) >= lines_per_page:
            chunk_text = "\n".join(buf)
            est = base_line // lines_per_page + 1
            label = f"约第{est}页(按行估算, 行{base_line+1}-{i+1})"
            out.extend(
                chunk_unit_text(
                    chunk_text,
                    source=name,
                    page_label=label,
                    meta="文本",
                    ext=".txt",
                )
            )
            buf = []
            base_line = i + 1
    if buf:
        est = base_line // lines_per_page + 1
        chunk_text = "\n".join(buf)
        label = f"约第{est}页(按行估算, 行{base_line+1}-{len(lines)})"
        out.extend(
            chunk_unit_text(
                chunk_text,
                source=name,
                page_label=label,
                meta="文本",
                ext=".txt",
            )
        )
    if not out and text.strip():
        out.extend(
            chunk_unit_text(
                text.strip(),
                source=name,
                page_label="全文",
                meta="文本",
                ext=".txt",
            )
        )
    return out


_rapid_ocr_engine: Any = None


def _get_rapid_ocr() -> Any:
    global _rapid_ocr_engine
    if _rapid_ocr_engine is None:
        py = sys.executable
        try:
            import onnxruntime  # noqa: F401  # RapidOCR 依赖；缺省时常表现为 rapidocr 导入失败
        except ImportError as e:
            raise ValueError(
                f"图像 OCR 依赖 onnxruntime。请在运行服务的同一 Python 环境中执行："
                f'"{py}" -m pip install onnxruntime rapidocr-onnxruntime pillow'
                f"\n（若使用 conda：先 conda activate forrag 再安装。）\n原始错误: {e!r}"
            ) from e
        try:
            from rapidocr_onnxruntime import RapidOCR
        except ImportError as e:
            raise ValueError(
                f"无法导入 rapidocr_onnxruntime。当前解释器: {py}\n"
                f'请执行: "{py}" -m pip install rapidocr-onnxruntime pillow\n'
                f"原始错误: {e!r}"
            ) from e
        try:
            _rapid_ocr_engine = RapidOCR()
        except Exception as e:
            raise RuntimeError(
                f"RapidOCR 初始化失败（解释器: {py}）。可尝试重装: "
                f'"{py}" -m pip install --force-reinstall onnxruntime rapidocr-onnxruntime\n'
                f"详情: {e!r}"
            ) from e
    return _rapid_ocr_engine


def _is_number_seq(obj: object) -> bool:
    return isinstance(obj, (list, tuple)) and bool(obj) and all(
        isinstance(x, (int, float)) for x in obj
    )


def _text_lines_from_rapidocr_output(ocr_out: object) -> list[str]:
    """从 RapidOCR 返回值中取出文本行。

    不同版本返回结构不同：(rows, 耗时) 的耗时可能是单个数字，也可能是各阶段耗时列表；
    新版本还可能返回带 .txts 属性的结果对象。
    """
    if ocr_out is None:
        return []
    txts = getattr(ocr_out, "txts", None)
    if txts is not None:
        return [str(t).strip() for t in txts if str(t).strip()]
    payload: object = ocr_out
    if isinstance(ocr_out, (list, tuple)) and len(ocr_out) == 2:
        a, b = ocr_out[0], ocr_out[1]
        if isinstance(b, (int, float)) or _is_number_seq(b) or b is None:
            payload = a
    rows = payload
    if not isinstance(rows, (list, tuple)):
        return []
    texts: list[str] = []
    for item in rows:
        if item is None:
            continue
        if isinstance(item, (list, tuple)) and len(item) >= 2:
            seg = item[1]
            if isinstance(seg, (list, tuple)) and len(seg) >= 1:
                texts.append(str(seg[0]).strip())
            elif isinstance(seg, str):
                texts.append(seg.strip())
        elif isinstance(item, str):
            texts.append(item.strip())
    return [t for t in texts if t]


def _strict_image_ocr() -> bool:
    """为真时：缺少 OCR 依赖直接报错；否则写入占位文本，避免整次问答失败。"""
    return os.environ.get("RAG_STRICT_IMAGE_OCR", "").strip().lower() in {"1", "true", "yes"}


def _parse_image_placeholder_chunks(name: str, detail: str) -> List[TextChunk]:
    """无 RapidOCR 时仍生成可检索占位块（提示安装依赖）。"""
    text = (
        f"【图片文件】{name}\n"
        f"（未能进行 OCR：{detail}\n"
        f"当前 Python：{sys.executable}\n"
        "请在本环境执行：python -m pip install -r requirements.txt\n"
        "若使用 conda：先 conda activate <你的环境名> 再安装并启动服务。）"
    )
    return chunk_unit_text(
        text,
        source=name,
        page_label="图片",
        meta="图像(无OCR)",
        ext=".png",
    )


def parse_image(path: Path) -> List[TextChunk]:
    """使用 RapidOCR 将图片中的文字识别为文本后分块；依赖缺失时可降级为占位文本。"""
    try:
        from PIL import Image
    except ImportError as e:
        raise ValueError("图像解析需要 pillow：pip install pillow") from e

    p = path.expanduser().resolve()
    name = p.name
    img = Image.open(p)
    if getattr(img, "n_frames", 1) > 1:
        img.seek(0)
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")

    strict = _strict_image_ocr()
    try:
        ocr = _get_rapid_ocr()
    except (ValueError, RuntimeError) as e:
        if strict:
            raise
        return _parse_image_placeholder_chunks(name, str(e))

    tmp_path: Optional[str] = None
    try:
        fd, tmp_path = tempfile.mkstemp(suffix=".png")
        os.close(fd)
        img.save(tmp_path, format="PNG")
        ocr_raw = ocr(tmp_path)
    except Exception as e:
        if strict:
            raise
        return _parse_image_placeholder_chunks(name, repr(e))
    finally:
        if tmp_path:
            Path(tmp_path).unlink(missing_ok=True)

    lines = _text_lines_from_rapidocr_output(ocr_raw)
    text = "\n".join(lines).strip()
    if not text:
        text = "（图片中未识别到文字内容）"

    return chunk_unit_text(
        text,
        source=name,
        page_label="图片",
        meta="图像 OCR",
        ext=path.suffix.lower(),
    )


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".pptx": parse_pptx,
    ".csv": parse_csv,
    ".txt": parse_txt,
    ".md": parse_txt,
}

_IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff")
for _ext in _IMAGE_EXTS:
    PARSERS[_ext] = parse_image


def parse_document(path: Path) -> List[TextChunk]:
    p = path.expanduser().resolve()
    ext = p.suffix.lower()
    fn = PARSERS.get(ext)
    if not fn:
        raise ValueError(f"不支持的扩展名 {ext}")
    raw = fn(p)
    cleaned = _dedupe_chunks(_merge_short_chunks(raw, ext))
    chunks = _finalize_chunks(cleaned, p)
    dropped = len(raw) - len(chunks)
    suffix = f"（合并/去重 {dropped} 个短块或重复块）" if dropped > 0 else ""
    print(f"[解析] {p.name} → {len(chunks)} 个文本块{suffix}")
    return chunks


def load_documents(paths: Sequence[Path]) -> List[TextChunk]:
    all_chunks: List[TextChunk] = []
    for p in paths:
        p = p.expanduser().resolve()
        if not p.is_file():
            print(f"[跳过] 不存在: {p}", file=sys.stderr)
            continue
        try:
            all_chunks.extend(parse_document(p))
        except Exception as e:
            print(f"[错误] 解析失败 {p}: {e}", file=sys.stderr)
    return all_chunks


# ------------- 向量检索 -------------


def _log_step(msg: str) -> None:
    print(msg, flush=True)


def _cache_root() -> Path:
    env = os.environ.get("RAG_CACHE_ROOT", "").strip()
    if env:
        root = Path(env).expanduser().resolve()
        root.mkdir(parents=True, exist_ok=True)
        return root
    return Path(__file__).resolve().parent / ".rag_cache"


def invalidate_caches_for_file(path: Path, embed_model_id: str) -> None:
    """
    删除指定文件对应的解析缓存与单文档向量缓存，并移除引用该文件路径的整库 bundle 缓存。
    必须在物理文件仍存在时调用（单文档缓存 key 依赖文件指纹）。
    """
    p = Path(path).expanduser().resolve()
    if not p.is_file():
        return
    norm = _normalized_path(p)
    doc = _doc_cache_paths(p, embed_model_id)
    shutil.rmtree(doc["root"], ignore_errors=True)
    shutil.rmtree(_parse_cache_paths(p)["root"], ignore_errors=True)

    bundles_dir = _cache_root() / "bundles"
    if not bundles_dir.is_dir():
        return
    for child in list(bundles_dir.iterdir()):
        if not child.is_dir():
            continue
        man = child / "manifest.json"
        if not man.is_file():
            continue
        try:
            data = json.loads(man.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            continue
        for item in data.get("files") or []:
            if item.get("path") == norm:
                shutil.rmtree(child, ignore_errors=True)
                break


def _bundle_key(paths: Sequence[Path], embed_model_id: str, bundle_extra: str = "") -> str:
    payload: dict[str, Any] = {
        "version": CACHE_VERSION,
        "embed_model_id": embed_model_id,
        "files": [
            {"path": _normalized_path(p), "fingerprint": _file_fingerprint(Path(p))}
            for p in sorted((Path(p) for p in paths), key=lambda x: str(x))
        ],
    }
    if bundle_extra:
        payload["kb_bundle"] = bundle_extra
    return _sha1_text(json.dumps(payload, ensure_ascii=False, sort_keys=True))


def _doc_cache_paths(path: Path, embed_model_id: str) -> dict:
    key = _sha1_text(
        json.dumps(
            {
                "version": CACHE_VERSION,
                "embed_model_id": embed_model_id,
                "file": _normalized_path(path),
                "fingerprint": _file_fingerprint(path),
            },
            ensure_ascii=False,
            sort_keys=True,
        )
    )
    root = _cache_root() / "docs" / key
    return {
        "root": root,
        "manifest": root / "manifest.json",
        "chunks": root / "chunks.jsonl",
        "embeddings": root / "embeddings.npy",
    }


def _parse_cache_paths(path: Path) -> dict:
    """解析结果缓存：只依赖分块版本与文件指纹，**不含嵌入模型**。

    解析是这里最贵的一步（图片页 OCR 每页数秒），换嵌入模型或做检索 A/B 时不应重跑。
    """
    key = _sha1_text(
        json.dumps(
            {
                "version": CACHE_VERSION,
                "file": _normalized_path(path),
                "fingerprint": _file_fingerprint(path),
            },
            ensure_ascii=False,
            sort_keys=True,
        )
    )
    root = _cache_root() / "parsed" / key
    return {"root": root, "manifest": root / "manifest.json", "chunks": root / "chunks.jsonl"}


def parse_document_cached(path: Path) -> List[TextChunk]:
    """带缓存的解析；缓存损坏时静默重解析。"""
    files = _parse_cache_paths(path)
    if files["manifest"].is_file() and files["chunks"].is_file():
        try:
            return _read_chunks(files["chunks"])
        except (OSError, json.JSONDecodeError, KeyError):
            shutil.rmtree(files["root"], ignore_errors=True)
    chunks = parse_document(path)
    try:
        _write_chunks(files["chunks"], chunks)
        files["manifest"].write_text(
            json.dumps(
                {
                    "version": CACHE_VERSION,
                    "file_path": _normalized_path(path),
                    "file_fingerprint": _file_fingerprint(path),
                    "chunk_count": len(chunks),
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
    except OSError:
        shutil.rmtree(files["root"], ignore_errors=True)
    return chunks


def _bundle_cache_paths(paths: Sequence[Path], embed_model_id: str, bundle_extra: str = "") -> dict:
    key = _bundle_key(paths, embed_model_id, bundle_extra)
    root = _cache_root() / "bundles" / key
    return {
        "root": root,
        "manifest": root / "manifest.json",
        "chunks": root / "chunks.jsonl",
        "embeddings": root / "embeddings.npy",
        "index": root / "faiss.index",
    }


def _write_chunks(path: Path, chunks: Sequence[TextChunk]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        for chunk in chunks:
            f.write(json.dumps(_chunk_to_dict(chunk), ensure_ascii=False) + "\n")


def _read_chunks(path: Path) -> List[TextChunk]:
    chunks: List[TextChunk] = []
    with path.open("r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            chunks.append(_chunk_from_dict(json.loads(line)))
    return chunks


def _embed_prefixes(model_id: str) -> tuple[str, str]:
    """返回 (query_prefix, passage_prefix)。

    e5 系列（multilingual-e5 等）要求查询/段落分别加 "query: " / "passage: " 前缀，
    否则检索质量明显下降。其它模型（bge、gte…）默认不加前缀，保持既有行为。
    可用 RAG_EMBED_QUERY_PREFIX / RAG_EMBED_PASSAGE_PREFIX 显式覆盖。
    """
    q = os.environ.get("RAG_EMBED_QUERY_PREFIX")
    p = os.environ.get("RAG_EMBED_PASSAGE_PREFIX")
    if q is not None or p is not None:
        return (q or "", p or "")
    if "e5" in (model_id or "").lower():
        return ("query: ", "passage: ")
    return ("", "")


def _load_st_model(embed_model_id: str):
    """同一 embed_model_id 全进程只加载一次；避免多线程并发拉 Hub 触发 httpx「client has been closed」等错误。"""
    with _st_model_singleton_lock:
        if embed_model_id in _st_model_singleton:
            return _st_model_singleton[embed_model_id]
        from sentence_transformers import SentenceTransformer

        local_path = os.environ.get("RAG_EMBED_MODEL_PATH", "").strip()
        model_ref = local_path or embed_model_id
        _log_step(
            f"[嵌入] (1/2) 加载向量模型 {model_ref}（首次会下载约百兆；默认 HF_ENDPOINT=镜像；也可先下载到本地并设置 RAG_EMBED_MODEL_PATH）…"
        )
        try:
            st = SentenceTransformer(model_ref)
        except Exception:
            _log_step(
                "[嵌入] 加载失败。若走 Hub：检查网络/代理，或显式 set HF_ENDPOINT=… 后重启；需官方源时 set RAG_USE_OFFICIAL_HF=1。若走本地：检查 RAG_EMBED_MODEL_PATH 是否指向完整模型目录。"
            )
            raise
        qp, pp = _embed_prefixes(embed_model_id)
        st._rag_query_prefix = qp
        st._rag_passage_prefix = pp
        _st_model_singleton[embed_model_id] = st
        return st


def _normalize_doc_embeddings(doc_emb: np.ndarray, st, num_chunks: int) -> np.ndarray:
    """保证每份文档的嵌入为 (n, dim)；0 条块时 sentence-transformers 常返回 1D 空数组，不能直接 concatenate。"""
    doc_emb = np.asarray(doc_emb, dtype=np.float32)
    dim = int(st.get_sentence_embedding_dimension())
    if num_chunks == 0:
        return np.zeros((0, dim), dtype=np.float32)
    if doc_emb.ndim == 1:
        return doc_emb.reshape(1, -1)
    return doc_emb


def chunk_embed_text(chunk: TextChunk) -> str:
    """用于检索的块文本：Contextual Retrieval 语境头 + 正文（正文本身仍用于 prompt/展示）。"""
    header = getattr(chunk, "context_header", "") or ""
    body = chunk.text or ""
    return f"{header}\n{body}" if header else body


def _embed_batch_size() -> int:
    """编码批大小；8GB CPU 机上默认 4，避免 sentence-transformers 默认 32 把内存打爆。"""
    return _env_int("RAG_EMBED_BATCH_SIZE", 4, minimum=1)


def _encode_chunks(st, chunks: Sequence[TextChunk]) -> np.ndarray:
    passage_prefix = getattr(st, "_rag_passage_prefix", "") or ""
    texts = [passage_prefix + chunk_embed_text(c) for c in chunks]
    batch = _embed_batch_size()
    _log_step(f"[嵌入] (2/2) 编码 {len(texts)} 条文本块（batch={batch}）…")
    if not texts:
        dim = int(st.get_sentence_embedding_dimension())
        return np.zeros((0, dim), dtype=np.float32)
    emb = st.encode(
        list(texts),
        batch_size=batch,
        show_progress_bar=True,
        normalize_embeddings=True,
    )
    out = np.asarray(emb, dtype=np.float32)
    if out.ndim == 1:
        out = out.reshape(1, -1)
    return out


def _build_faiss_index(embeddings: np.ndarray):
    import faiss

    index = faiss.IndexFlatIP(embeddings.shape[1])
    index.add(embeddings)
    return index


def _save_doc_cache(path: Path, embed_model_id: str, chunks: Sequence[TextChunk], embeddings: np.ndarray) -> None:
    files = _doc_cache_paths(path, embed_model_id)
    files["root"].mkdir(parents=True, exist_ok=True)
    _write_chunks(files["chunks"], chunks)
    np.save(files["embeddings"], embeddings)
    manifest = {
        "version": CACHE_VERSION,
        "embed_model_id": embed_model_id,
        "file_path": _normalized_path(path),
        "file_fingerprint": _file_fingerprint(path),
        "chunk_count": len(chunks),
    }
    files["manifest"].write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")


def _load_doc_cache(path: Path, embed_model_id: str) -> tuple[List[TextChunk], np.ndarray] | None:
    files = _doc_cache_paths(path, embed_model_id)
    if not files["manifest"].is_file() or not files["chunks"].is_file() or not files["embeddings"].is_file():
        return None
    chunks = _read_chunks(files["chunks"])
    embeddings = np.load(files["embeddings"])
    return chunks, np.asarray(embeddings, dtype=np.float32)


def _save_bundle_cache(
    paths: Sequence[Path],
    embed_model_id: str,
    chunks: Sequence[TextChunk],
    embeddings: np.ndarray,
    index,
    bundle_extra: str = "",
) -> None:
    import faiss

    files = _bundle_cache_paths(paths, embed_model_id, bundle_extra)
    files["root"].mkdir(parents=True, exist_ok=True)
    _write_chunks(files["chunks"], chunks)
    np.save(files["embeddings"], embeddings)
    _faiss_write_index(index, files["index"])
    manifest = {
        "version": CACHE_VERSION,
        "embed_model_id": embed_model_id,
        "bundle_key": _bundle_key(paths, embed_model_id, bundle_extra),
        "files": [
            {"path": _normalized_path(p), "fingerprint": _file_fingerprint(Path(p))}
            for p in sorted((Path(p) for p in paths), key=lambda x: str(x))
        ],
        "chunk_count": len(chunks),
    }
    files["manifest"].write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")


def _load_bundle_cache(paths: Sequence[Path], embed_model_id: str, bundle_extra: str = ""):
    import faiss

    files = _bundle_cache_paths(paths, embed_model_id, bundle_extra)
    if not files["manifest"].is_file() or not files["chunks"].is_file() or not files["embeddings"].is_file() or not files["index"].is_file():
        return None
    chunks = _read_chunks(files["chunks"])
    embeddings = np.load(files["embeddings"])
    index = _faiss_read_index(files["index"])
    return chunks, np.asarray(embeddings, dtype=np.float32), index


def _apply_path_chunk_tags(doc_chunks: Sequence[TextChunk], path: Path, tags: Optional[dict[str, str]]) -> None:
    if not tags:
        return
    norm = _normalized_path(path)
    _ = norm  # reserved for future per-path diagnostics
    for c in doc_chunks:
        if tags.get("kb_note_id"):
            c.kb_note_id = tags["kb_note_id"]
        if tags.get("kb_attachment_id"):
            c.kb_attachment_id = tags["kb_attachment_id"]
        if tags.get("session_file_id"):
            c.session_file_id = tags["session_file_id"]


def build_or_load_index(
    paths: Sequence[Path],
    embed_model_id: str,
    bundle_extra: str = "",
    chunk_tags_by_norm_path: Optional[dict[str, dict[str, str]]] = None,
):
    cached_bundle = _load_bundle_cache(paths, embed_model_id, bundle_extra)
    st = _load_st_model(embed_model_id)
    if cached_bundle is not None:
        chunks, embeddings, index = cached_bundle
        _log_step(f"[缓存] 命中整库缓存，直接复用 {len(chunks)} 个文本块。")
        return chunks, embeddings, index, st

    chunks_all: List[TextChunk] = []
    emb_all: List[np.ndarray] = []
    changed_count = 0
    for path in paths:
        p = Path(path).expanduser().resolve()
        if not p.is_file():
            print(f"[跳过] 不存在: {p}", file=sys.stderr)
            continue
        ext = p.suffix.lower()
        if ext not in PARSERS:
            print(f"[跳过] 不支持的扩展名 {ext}: {p}", file=sys.stderr)
            continue
        tags = (chunk_tags_by_norm_path or {}).get(_normalized_path(p))
        cached_doc = _load_doc_cache(p, embed_model_id)
        if cached_doc is not None:
            doc_chunks, doc_emb = cached_doc
            _log_step(f"[缓存] 命中文档缓存：{p.name}（{len(doc_chunks)} 块）")
        else:
            doc_chunks = parse_document_cached(p)
            doc_emb = _encode_chunks(st, doc_chunks)
            _save_doc_cache(p, embed_model_id, doc_chunks, doc_emb)
            changed_count += 1
        _apply_path_chunk_tags(doc_chunks, p, tags)
        chunks_all.extend(doc_chunks)
        emb_all.append(_normalize_doc_embeddings(doc_emb, st, len(doc_chunks)))

    if not chunks_all:
        raise ValueError("没有可用的文档块，请检查文件路径与格式。")
    embeddings = np.concatenate(emb_all, axis=0).astype(np.float32, copy=False)
    index = _build_faiss_index(embeddings)
    _save_bundle_cache(paths, embed_model_id, chunks_all, embeddings, index, bundle_extra)
    if changed_count:
        _log_step(f"[缓存] 已更新 {changed_count} 个文档缓存，并重建整库索引。")
    return chunks_all, embeddings, index, st


def search(
    query: str,
    chunks: Sequence[TextChunk],
    index,
    st,
    top_k: int,
) -> List[tuple[float, TextChunk]]:
    query_prefix = getattr(st, "_rag_query_prefix", "") or ""
    q = st.encode([query_prefix + query], normalize_embeddings=True)
    q = np.asarray(q, dtype=np.float32)
    n = len(chunks)
    if n == 0:
        return []
    k = max(1, min(int(top_k), n))
    scores, ids = index.search(q, k)
    out: List[tuple[float, TextChunk]] = []
    for score, idx in zip(scores[0], ids[0]):
        if idx < 0:
            continue
        out.append((float(score), chunks[idx]))
    return out


def route_generation(has_api_key: bool) -> str:
    if has_api_key:
        return "api"
    return "local"


# ------------- 大模型（魔搭 / Hugging Face）-------------


def _configure_llm_loading() -> None:
    """减轻「像卡死」：打开 transformers 日志、限制线程，避免占满 CPU。"""
    import os
    import torch

    try:
        import transformers

        transformers.utils.logging.set_verbosity_info()
    except Exception:
        pass
    n = os.cpu_count() or 4
    # 加载大权重时少占线程，系统更不容易假死
    nt = max(1, min(4, n // 2))
    try:
        torch.set_num_threads(nt)
        torch.set_num_interop_threads(1)
    except Exception:
        pass


def _llm_dtype(cpu_half: bool = False):
    import torch

    if torch.cuda.is_available():
        return torch.bfloat16 if torch.cuda.is_bf16_supported() else torch.float16
    # 无 GPU：fp32 约需 12GB+ 内存（3B），易因换页长时间无响应。
    # 注意：float16 在很多 CPU 上反而会非常慢，看起来像“完全不出字”。
    if cpu_half:
        cpu_has_bf16 = bool(getattr(getattr(torch.backends, "cpu", None), "has_bf16", False))
        if cpu_has_bf16:
            _log_step("[LLM] 检测到 CPU 支持 bfloat16，使用 bfloat16 以降低内存占用。")
            return torch.bfloat16
        _log_step("[LLM] 警告：当前 CPU 不适合 float16 推理；为避免生成阶段极慢，改回 float32。")
    return torch.float32


def _normalize_llm_hub(name: str) -> str:
    n = (name or "auto").strip().lower()
    if n in ("ms", "modelscope"):
        return "modelscope"
    if n in ("hf", "huggingface"):
        return "huggingface"
    return "auto"


def _from_pretrained_llm(model_id: str, dtype, hub_ms: bool):
    """避免 CPU 上 device_map='auto' 与部分 transformers 版本组合时出现 meta 张量错误。"""
    import torch

    if hub_ms:
        from modelscope import AutoModelForCausalLM
    else:
        from transformers import AutoModelForCausalLM

    def load(**extra):
        return AutoModelForCausalLM.from_pretrained(model_id, trust_remote_code=True, **extra)

    if torch.cuda.is_available():
        _log_step("[LLM] 正在将权重载入 GPU（device_map=auto，请稍候）…")
        try:
            m = load(dtype=dtype, device_map="auto")
        except TypeError:
            m = load(torch_dtype=dtype, device_map="auto")
        _log_step("[LLM] 权重已载入 GPU。")
        return m

    # CPU：整模放内存，不用 accelerate 分片，避免 generate 里 eos/pad 与 meta 张量冲突
    _log_step(
        "[LLM] 正在从磁盘读取并映射权重到内存（CPU 上可能需数分钟，磁盘灯会闪，并非死机）…"
    )
    try:
        m = load(dtype=dtype, device_map=None, low_cpu_mem_usage=True)
    except TypeError:
        m = load(torch_dtype=dtype, device_map=None, low_cpu_mem_usage=True)
    _log_step("[LLM] 权重加载完成，准备推理。")
    return m


def _fix_tokenizer_pad(tokenizer) -> None:
    if tokenizer.pad_token_id is None and tokenizer.eos_token_id is not None:
        tokenizer.pad_token = tokenizer.eos_token


def _sync_gen_config(model, tokenizer) -> None:
    gc = getattr(model, "generation_config", None)
    if gc is None:
        return
    if tokenizer.pad_token_id is not None:
        gc.pad_token_id = int(tokenizer.pad_token_id)
    if tokenizer.eos_token_id is not None:
        gc.eos_token_id = int(tokenizer.eos_token_id)
    # 与 generate(do_sample=False) 一致，避免与仓库里 temperature/top_p/top_k 冲突刷屏
    gc.do_sample = False
    for name in ("temperature", "top_p", "top_k"):
        if hasattr(gc, name):
            try:
                setattr(gc, name, None)
            except Exception:
                pass


def _looks_like_model_dir(p: Path) -> bool:
    return (p / "config.json").is_file() or (p / "tokenizer_config.json").is_file()


def _find_modelscope_local_snapshot(model_id: str) -> str | None:
    """魔搭把模型放在 ~/.cache/modelscope/hub/models/组织/名称/，名称里的点会变成 ___。"""
    base = Path(os.environ.get("MODELSCOPE_CACHE", Path.home() / ".cache" / "modelscope")) / "hub" / "models"
    if not base.is_dir() or "/" not in model_id:
        return None
    org, name = model_id.split("/", 1)
    candidates = [
        base / org / name,
        base / org / name.replace(".", "___"),
    ]
    for c in candidates:
        if _looks_like_model_dir(c):
            return str(c.resolve())
    org_dir = base / org
    if org_dir.is_dir():
        for child in sorted(org_dir.iterdir(), key=lambda x: x.name):
            if child.is_dir() and _looks_like_model_dir(child) and name.split("-")[-1] in child.name:
                return str(child.resolve())
    return None


def load_llm(model_id: str, hub: str = "auto", cpu_half: bool = False):
    """hub: auto | modelscope | huggingface。auto 在魔搭 DNS/网络失败时回退到 Hugging Face。"""
    _configure_llm_loading()
    dtype = _llm_dtype(cpu_half=cpu_half)
    hub_n = _normalize_llm_hub(hub)

    def from_hf():
        from transformers import AutoTokenizer

        _log_step(f"[LLM] 从 Hugging Face 加载: {model_id}（首次会下载权重）")
        _log_step("[LLM] (1/2) 加载 tokenizer …")
        tokenizer = AutoTokenizer.from_pretrained(model_id, trust_remote_code=True)
        _fix_tokenizer_pad(tokenizer)
        _log_step("[LLM] (2/2) 加载生成模型权重（最耗时）…")
        model = _from_pretrained_llm(model_id, dtype, hub_ms=False)
        model.eval()
        _sync_gen_config(model, tokenizer)
        return model, tokenizer

    def from_local_dir(local_dir: str, label: str = "本地目录"):
        from transformers import AutoTokenizer

        _log_step(f"[LLM] 从{label}加载（不访问魔搭在线接口）: {local_dir}")
        _log_step("[LLM] (1/2) 加载 tokenizer …")
        tokenizer = AutoTokenizer.from_pretrained(local_dir, trust_remote_code=True)
        _fix_tokenizer_pad(tokenizer)
        _log_step("[LLM] (2/2) 加载生成模型权重 …")
        model = _from_pretrained_llm(local_dir, dtype, hub_ms=False)
        model.eval()
        _sync_gen_config(model, tokenizer)
        return model, tokenizer

    def from_ms_online():
        from modelscope import AutoTokenizer

        _log_step(f"[LLM] 从魔搭 ModelScope 在线加载: {model_id}（首次会下载权重）")
        _log_step("[LLM] (1/2) 加载 tokenizer …")
        tokenizer = AutoTokenizer.from_pretrained(model_id, trust_remote_code=True)
        _fix_tokenizer_pad(tokenizer)
        _log_step("[LLM] (2/2) 加载生成模型权重 …")
        model = _from_pretrained_llm(model_id, dtype, hub_ms=True)
        model.eval()
        _sync_gen_config(model, tokenizer)
        return model, tokenizer

    def try_ms_with_local_fallback():
        try:
            return from_ms_online()
        except Exception as e:
            local = _find_modelscope_local_snapshot(model_id)
            if local:
                print(
                    f"[LLM] 魔搭在线校验失败（{type(e).__name__}），改用本机已下载缓存。",
                    file=sys.stderr,
                )
                return from_local_dir(local, label="魔搭缓存")
            raise

    if hub_n == "huggingface":
        return from_hf()
    if hub_n == "modelscope":
        try:
            return try_ms_with_local_fallback()
        except Exception as e:
            print(
                "[LLM] 魔搭不可用且无本地缓存。可改用: --llm-hub huggingface "
                "或设置 $env:HF_ENDPOINT='https://hf-mirror.com'",
                file=sys.stderr,
            )
            raise e

    try:
        return try_ms_with_local_fallback()
    except Exception as e:
        print(
            f"[LLM] 自动模式：魔搭/本地缓存均不可用（{type(e).__name__}），改用 Hugging Face。摘要: {e!s}",
            file=sys.stderr,
        )
        print(
            "[提示] 若 HuggingFace 较慢: $env:HF_ENDPOINT='https://hf-mirror.com'",
            file=sys.stderr,
        )
        return from_hf()


def build_prompt(
    question: str,
    retrieved: List[tuple[float, TextChunk]],
) -> str:
    parts = []
    for i, (score, ch) in enumerate(retrieved, 1):
        ref = f"[片段{i}] 文件: {ch.source} | 位置: {ch.page_label} | {ch.meta}"
        parts.append(f"{ref}\n内容: {ch.text}\n")
    context = "\n".join(parts)
    user_msg = (
        "你是文档问答助手。请仅根据下面「检索到的文档片段」回答问题。\n"
        "要求：\n"
        "1. 先给出简洁答案。\n"
        "2. 在答案中或末尾用括号标明依据：文件名 + 页码/位置（与片段中的「位置」一致）。\n"
        "3. 若片段不足以回答，请自由回答。\n\n"
        f"用户问题：{question}\n\n"
        f"检索到的文档片段：\n{context}"
    )
    return user_msg


def generate_answer(
    model,
    tokenizer,
    user_msg: str,
    max_new_tokens: Optional[int],
    stream: bool = True,
) -> str:
    import torch
    from transformers import TextStreamer

    messages = [{"role": "user", "content": user_msg}]
    # Qwen2.5 / Qwen3 均支持 chat_template
    try:
        text = tokenizer.apply_chat_template(
            messages,
            tokenize=False,
            add_generation_prompt=True,
        )
    except Exception:
        text = user_msg

    pad_id = tokenizer.pad_token_id
    if pad_id is None:
        pad_id = tokenizer.eos_token_id
    eos_id = tokenizer.eos_token_id

    inputs = tokenizer(text, return_tensors="pt")
    device = next(model.parameters()).device
    inputs = {k: v.to(device) for k, v in inputs.items()}
    gen_kw = {
        "do_sample": False,
        "pad_token_id": pad_id,
        "eos_token_id": eos_id,
    }
    if max_new_tokens is not None:
        gen_kw["max_new_tokens"] = max_new_tokens
    if not stream:
        with torch.no_grad():
            out = model.generate(**inputs, **gen_kw)
        new_tokens = out[0, inputs["input_ids"].shape[1] :]
        reply = tokenizer.decode(new_tokens, skip_special_tokens=True)
        return reply.strip()

    _log_step("[LLM] 开始流式生成回答：")
    streamer = TextStreamer(
        tokenizer,
        skip_prompt=True,
        skip_special_tokens=True,
    )
    with torch.no_grad():
        out = model.generate(**inputs, **gen_kw, streamer=streamer)
    print(flush=True)
    new_tokens = out[0, inputs["input_ids"].shape[1] :]
    reply = tokenizer.decode(new_tokens, skip_special_tokens=True)
    return reply.strip()


def generate_answer_via_api(
    api_key: str,
    api_model: str,
    api_base: str,
    user_msg: str,
    max_new_tokens: Optional[int],
    stream: bool = True,
    json_object: bool = False,
) -> str:
    import httpx
    from openai import OpenAI
    from openai import APIConnectionError

    timeout = httpx.Timeout(120.0, connect=20.0)
    with httpx.Client(http2=False, trust_env=False, timeout=timeout) as http_client:
        client = OpenAI(
            api_key=api_key,
            base_url=api_base,
            http_client=http_client,
            max_retries=2,
        )
        req: dict[str, Any] = {
            "model": api_model,
            "messages": [{"role": "user", "content": user_msg}],
            "temperature": 0,
            "stream": stream,
        }
        if max_new_tokens is not None:
            req["max_tokens"] = max_new_tokens
        if json_object and not stream:
            req["response_format"] = {"type": "json_object"}
        if not stream:
            try:
                resp = client.chat.completions.create(**req)
                return (resp.choices[0].message.content or "").strip()
            except Exception:
                if json_object and "response_format" in req:
                    req.pop("response_format", None)
                    resp = client.chat.completions.create(**req)
                    return (resp.choices[0].message.content or "").strip()
                raise

        _log_step(f"[API] 使用千问兼容 API 流式生成：{api_model}")
        chunks: List[str] = []
        try:
            resp = client.chat.completions.create(**req)
            for event in resp:
                delta = event.choices[0].delta.content or ""
                if not delta:
                    continue
                chunks.append(delta)
                print(delta, end="", flush=True)
            print(flush=True)
            return "".join(chunks).strip()
        except APIConnectionError:
            _log_step("[API] 流式连接失败，自动退回非流式请求重试…")
            req["stream"] = False
            resp = client.chat.completions.create(**req)
            text = (resp.choices[0].message.content or "").strip()
            print(text, flush=True)
            return text


def main():
    parser = argparse.ArgumentParser(description="魔搭 Qwen 文档问答（上传办公文件，检索页码与答案）")
    parser.add_argument("--files", nargs="+", required=True, help="一个或多个文件路径")
    parser.add_argument("--question", "-q", required=True, help="用户问题")
    parser.add_argument(
        "--top-k",
        type=int,
        default=3,
        help="向量检索返回的片段条数，默认 3（仅影响送入模型的上下文，不改变索引缓存）",
    )
    parser.add_argument(
        "--max-new-tokens",
        type=int,
        default=None,
        help="生成新 token 上限；不设则不在请求里限制长度（由服务端/模型自行决定）",
    )
    parser.add_argument(
        "--model",
        default=os.environ.get("MS_MODEL_ID", "Qwen/Qwen2.5-3B-Instruct"),
        help="模型 ID（魔搭与 HF 上 Qwen 常同名），环境变量 MS_MODEL_ID；例：Qwen/Qwen2.5-7B-Instruct",
    )
    parser.add_argument(
        "--api-key",
        default=os.environ.get("DASHSCOPE_API_KEY", ""),
        help="千问 API Key；提供后将直接走 API，不加载本地大模型。也可设环境变量 DASHSCOPE_API_KEY",
    )
    parser.add_argument(
        "--api-model",
        default=os.environ.get("QWEN_API_MODEL", "qwen-plus"),
        help="千问 API 模型名，默认 qwen-plus；也可设环境变量 QWEN_API_MODEL",
    )
    parser.add_argument(
        "--api-base",
        default=os.environ.get(
            "QWEN_API_BASE", "https://dashscope.aliyuncs.com/compatible-mode/v1"
        ),
        help="千问兼容 API Base URL；也可设环境变量 QWEN_API_BASE",
    )
    parser.add_argument(
        "--llm-hub",
        default=_normalize_llm_hub(os.environ.get("LLM_HUB", "auto")),
        choices=["auto", "modelscope", "huggingface"],
        help="auto=魔搭在线→失败则用本机魔搭缓存→再失败则 HF；modelscope=优先在线，失败用缓存；huggingface=仅 HF。环境变量 LLM_HUB",
    )
    parser.add_argument(
        "--embed-model",
        default=os.environ.get("MS_EMBED_ID", "BAAI/bge-small-zh-v1.5"),
        help="向量检索嵌入模型 ID（默认 HuggingFace 的 bge-small-zh；也可换魔搭上的中文句向量模型并设 MS_EMBED_ID）",
    )
    parser.add_argument(
        "--low-memory",
        action="store_true",
        help="无 GPU 时用 float16 加载大模型，内存约减半（推荐 16GB 以下内存或加载卡住时尝试）",
    )
    parser.add_argument(
        "--no-stream",
        action="store_true",
        help="关闭实时 token 输出，改为生成完成后一次性打印",
    )
    args = parser.parse_args()

    paths = [Path(f) for f in args.files]
    try:
        chunks, _embeddings, index, st = build_or_load_index(paths, args.embed_model)
    except ValueError as e:
        print(str(e), file=sys.stderr)
        sys.exit(1)

    hits = search(args.question, chunks, index, st, args.top_k)

    print("\n---------- 检索到的位置（Top-%d）----------" % len(hits))
    for score, ch in hits:
        print(f"  相关度 {score:.4f} | {ch.source} | {ch.page_label} | {ch.meta}")

    print("\n---------- 模型生成回答 ----------")
    route = route_generation(has_api_key=bool(args.api_key.strip()))

    user_msg = build_prompt(args.question, hits)
    if route == "api":
        answer = generate_answer_via_api(
            api_key=args.api_key.strip(),
            api_model=args.api_model,
            api_base=args.api_base,
            user_msg=user_msg,
            max_new_tokens=args.max_new_tokens,
            stream=not args.no_stream,
        )
        if args.no_stream:
            print(answer)
        return

    _log_step("[路由] 当前问题走本地模型生成。")
    model, tokenizer = load_llm(args.model, hub=args.llm_hub, cpu_half=args.low_memory)
    answer = generate_answer(
        model, tokenizer, user_msg, args.max_new_tokens, stream=not args.no_stream
    )
    if args.no_stream:
        print(answer)


if __name__ == "__main__":
    main()
