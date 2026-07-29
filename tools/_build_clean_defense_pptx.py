#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Clean final-defense PPT (no HKU.pptx template).
Palette from user swatches; HKU wordmark top-right (image 2).
Technical content follows the dissertation PDF.
"""
from __future__ import annotations

import io
import time
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- paths ---
OUT_DIR = Path(r"C:\Users\10603\Desktop\毕业论文模版")
ASSETS = Path(r"C:\Users\10603\.cursor\projects\e-For-RAG\assets")
LOGO = ASSETS / "c__Users_10603_AppData_Roaming_Cursor_User_workspaceStorage_210eba0554e21ef301030bf7636ca08f_images_image-ca45a952-8e13-4cad-972b-165764d29b98.png"
FIGS = Path(r"C:\Users\10603\Desktop\毕业论文模版\第十一稿\figures")
OUT_NAME = "Jin_Yue_3036494489_Final_Defense_Clean.pptx"

# --- palette (from user images) ---
GREEN_DARK = RGBColor(0x27, 0x58, 0x29)      # image 1 / 4
GREEN_DEEPER = RGBColor(0x1B, 0x37, 0x06)
GREEN_LIGHT = RGBColor(0x61, 0xBB, 0x45)     # image 3
GREEN_MID = RGBColor(0x56, 0xA8, 0x40)       # image 5
TEXT = RGBColor(0x1F, 0x2A, 0x2E)
MUTED = RGBColor(0x5A, 0x6B, 0x65)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF9, 0xF7)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD7, 0xE2, 0xD8)

# Typography
FONT = "Times New Roman"
# Size scale tuned for 16:9 + Times New Roman
SZ_COVER = 26          # cover main title
SZ_COVER_SUB = 14      # cover meta
SZ_SLIDE_TITLE = 24    # section title in header
SZ_CARD_TITLE = 14     # card / block titles
SZ_BODY = 12           # body bullets
SZ_LEAD = 13           # intro / lead sentence
SZ_LABEL = 13          # layer tags, small headings
SZ_METRIC = 26         # big numbers
SZ_FOOTER = 10
SZ_THANK = 40

# 16:9
W, H = Inches(13.333), Inches(7.5)


def _rgb(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _apply_font(run, *, size, bold, color, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _set_text(tf, text: str, *, size=SZ_BODY, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font=FONT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _apply_font(run, size=size, bold=bold, color=color, font=font)


def _add_paras(
    tf,
    lines: list[str],
    *,
    size=SZ_BODY,
    color=TEXT,
    bold_first=False,
    spacing=6,
    font=FONT,
    sub_size=None,
):
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = line
        sz = size if i == 0 else (sub_size or size)
        _apply_font(run, size=sz, bold=(bold_first and i == 0), color=color, font=font)


def add_logo(slide):
    if LOGO.exists():
        # top-right, keep aspect ~433x85
        slide.shapes.add_picture(str(LOGO), Inches(9.55), Inches(0.22), width=Inches(3.4))


def add_footer(slide, page: int, total: int, section: str = ""):
    # bottom accent line
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.22), W, Inches(0.28))
    _rgb(bar, GREEN_DARK)
    # page number
    box = slide.shapes.add_textbox(Inches(11.6), Inches(7.24), Inches(1.5), Inches(0.24))
    _set_text(box.text_frame, f"{page} / {total}", size=SZ_FOOTER, color=WHITE, align=PP_ALIGN.RIGHT)
    if section:
        sb = slide.shapes.add_textbox(Inches(0.4), Inches(7.24), Inches(8), Inches(0.24))
        _set_text(sb.text_frame, section, size=SZ_FOOTER, color=WHITE)


def add_title_bar(slide, title: str):
    # left accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), H)
    _rgb(accent, GREEN_DARK)
    # header strip
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.12), Inches(0), Inches(13.213), Inches(0.95))
    _rgb(head, BG)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.12), Inches(0.94), Inches(13.213), Inches(0.03))
    _rgb(line, GREEN_LIGHT)
    # keep clear of top-right logo (~3.4" wide from 9.55)
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(8.7), Inches(0.62))
    tb.text_frame.word_wrap = True
    # slightly smaller when title is long
    sz = SZ_SLIDE_TITLE if len(title) < 42 else 20
    _set_text(tb.text_frame, title, size=sz, bold=True, color=GREEN_DARK)
    add_logo(slide)


def add_card(slide, left, top, width, height, title: str, body_lines: list[str], *, accent=GREEN_DARK):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    # top accent
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
    _rgb(strip, accent)
    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.18), width - Inches(0.36), Inches(0.42))
    title_box.text_frame.word_wrap = True
    _set_text(title_box.text_frame, title, size=SZ_CARD_TITLE, bold=True, color=GREEN_DARK)
    body = slide.shapes.add_textbox(
        left + Inches(0.18),
        top + Inches(0.62),
        width - Inches(0.36),
        height - Inches(0.78),
    )
    body.text_frame.word_wrap = True
    bullets = [f"• {x}" if not x.startswith("•") else x for x in body_lines]
    _add_paras(body.text_frame, bullets, size=SZ_BODY, color=TEXT, spacing=6)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides_meta = []  # (slide, section) for footer later
    total = 16  # planned

    # ========== 1 Title ==========
    s = new_slide(prs)
    # left panel
    left = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.6), H)
    _rgb(left, GREEN_DARK)
    # light accent block
    accent = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(5.4), Inches(3.6), Inches(1.2))
    _rgb(accent, GREEN_LIGHT)
    t1 = s.shapes.add_textbox(Inches(0.65), Inches(5.55), Inches(3.2), Inches(0.9))
    _add_paras(
        t1.text_frame,
        ["Final Oral Defense", "July 2026"],
        size=15,
        color=WHITE,
        bold_first=True,
        spacing=4,
        sub_size=13,
    )
    # right content
    add_logo(s)
    title = s.shapes.add_textbox(Inches(5.1), Inches(2.0), Inches(7.7), Inches(2.2))
    _add_paras(
        title.text_frame,
        [
            "Design and Implementation of a",
            "RAG-Based Teacher–Student",
            "Co-Learning System",
        ],
        size=SZ_COVER,
        color=GREEN_DARK,
        bold_first=False,
        spacing=2,
    )
    # make all title lines bold
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            _apply_font(r, size=SZ_COVER, bold=True, color=GREEN_DARK)

    meta = s.shapes.add_textbox(Inches(5.1), Inches(4.5), Inches(7.5), Inches(1.8))
    _add_paras(
        meta.text_frame,
        [
            "Jin Yue  ·  3036494489",
            "MSc(Eng) Electrical and Electronic Engineering",
            "The University of Hong Kong",
            "Supervisor: Dr. Andrew H.C. Wu",
        ],
        size=SZ_COVER_SUB,
        color=MUTED,
        spacing=6,
    )
    slides_meta.append((s, ""))

    # ========== 2 Overview ==========
    s = new_slide(prs)
    add_title_bar(s, "Overview")
    items = [
        ("01", "Task Scope", "What the system covers"),
        ("02", "Motivation", "Why this problem matters"),
        ("03", "Background Review", "RAG and related methods"),
        ("04", "System & Pipeline", "Technical implementation"),
        ("05", "Benchmark", "ELEC6081 evaluation set"),
        ("06", "Results", "Gate & re-ranking evidence"),
        ("07", "Visual Results", "Architecture figures"),
        ("08", "Conclusions & Future", "Takeaways and limits"),
    ]
    for i, (num, title, sub) in enumerate(items):
        col = i % 4
        row = i // 4
        left = Inches(0.45 + col * 3.15)
        top = Inches(1.35 + row * 2.6)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.95), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = LINE
        nbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), top + Inches(0.25), Inches(0.7), Inches(0.45))
        _rgb(nbox, GREEN_LIGHT if row == 0 else GREEN_DARK)
        nt = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.28), Inches(0.7), Inches(0.4))
        _set_text(nt.text_frame, num, size=SZ_LABEL, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.95), Inches(2.55), Inches(1.0))
        _add_paras(tb.text_frame, [title, sub], size=SZ_CARD_TITLE, color=TEXT, bold_first=True, spacing=6, sub_size=11)
        for j, p in enumerate(tb.text_frame.paragraphs):
            for r in p.runs:
                if j == 1:
                    _apply_font(r, size=11, bold=False, color=MUTED)
    slides_meta.append((s, "Overview"))

    # ========== 3 Task Scope ==========
    s = new_slide(prs)
    add_title_bar(s, "Task Scope")
    intro = s.shapes.add_textbox(Inches(0.45), Inches(1.15), Inches(12.4), Inches(0.45))
    intro.text_frame.word_wrap = True
    _set_text(
        intro.text_frame,
        "Goal: a course-deployable Teacher–Student Co-Learning (SOLO) Bot using retrieval-augmented generation.",
        size=SZ_LEAD,
        color=MUTED,
    )
    # three equal cards; leave clear gap above footer (7.22)
    card_top, card_h, card_w, gap = Inches(1.75), Inches(5.1), Inches(4.05), Inches(0.2)
    add_card(
        s,
        Inches(0.4),
        card_top,
        card_w,
        card_h,
        "Shared course knowledge base",
        [
            "One KB per course: teachers write, students read",
            "Formats: PDF, PPT, Word, Excel, images (+ OCR)",
            "OCR fallback for scanned / image pages",
            "Session uploads stay private to the session",
        ],
        accent=GREEN_DARK,
    )
    add_card(
        s,
        Inches(0.4) + card_w + gap,
        card_top,
        card_w,
        card_h,
        "Traceable RAG question answering",
        [
            "Hybrid retrieval: dense + BM25 + RRF",
            "Cross-encoder re-ranking with fallback",
            "Three-band gate: none / weak / grounded",
            "Citations target file + page/slide when metadata allows",
        ],
        accent=GREEN_LIGHT,
    )
    add_card(
        s,
        Inches(0.4) + 2 * (card_w + gap),
        card_top,
        card_w,
        card_h,
        "SOLO practice & feedback loop",
        [
            "Tiered exercises guided by SOLO taxonomy",
            "Automated grading against answer keys",
            "Teacher export of questions and scores",
            "Export omits full conversation transcripts",
        ],
        accent=GREEN_MID,
    )
    slides_meta.append((s, "Task Scope"))

    # ========== 4 Motivation ==========
    s = new_slide(prs)
    add_title_bar(s, "Motivation")
    problems = [
        ("Private course knowledge", "General LLMs are not kept in sync with latest slides, handouts, and course-specific materials."),
        ("Missing provenance", "Fluent answers without file/page evidence are hard to trust in teaching."),
        ("Incomplete learning loop", "QA alone does not provide practice, grading, and teacher feedback."),
        ("Messy real corpora", "PDFs, slides, tables, OCR pages, and bilingual terms often break naive pipelines."),
    )
    for i, (t, b) in enumerate(problems):
        top = Inches(1.25 + i * 1.35)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), top + Inches(0.15), Inches(0.55), Inches(0.55))
        _rgb(num, GREEN_DARK if i % 2 == 0 else GREEN_LIGHT)
        nt = s.shapes.add_textbox(Inches(0.55), top + Inches(0.22), Inches(0.55), Inches(0.45))
        _set_text(nt.text_frame, str(i + 1), size=SZ_LABEL, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), top, Inches(11.4), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = LINE
        tb = s.shapes.add_textbox(Inches(1.6), top + Inches(0.18), Inches(10.9), Inches(0.85))
        _add_paras(tb.text_frame, [t, b], size=SZ_CARD_TITLE, color=TEXT, bold_first=True, spacing=4, sub_size=SZ_BODY)
        for j, p in enumerate(tb.text_frame.paragraphs):
            for r in p.runs:
                if j == 1:
                    _apply_font(r, size=SZ_BODY, bold=False, color=MUTED)
    slides_meta.append((s, "Motivation"))

    # ========== 5 Background ==========
    s = new_slide(prs)
    add_title_bar(s, "Background Review")
    add_card(
        s,
        Inches(0.45),
        Inches(1.25),
        Inches(6.1),
        Inches(2.6),
        "RAG paradigm",
        [
            "Retriever + generator: condition answers on retrieved passages (Lewis et al.)",
            "Updates the non-parametric store without retraining the LLM",
            "Relative to direct answering: better material sync and traceability; can reduce unsupported claims",
        ],
        accent=GREEN_DARK,
    )
    add_card(
        s,
        Inches(6.8),
        Inches(1.25),
        Inches(6.05),
        Inches(2.6),
        "Faithfulness constraints",
        [
            "Reflection / corrective retrieval (e.g. CRAG-style three-band structure)",
            "Evidence position matters (Lost-in-the-Middle)",
            "Citation checks should be enforced, not left to prompting alone",
        ],
        accent=GREEN_LIGHT,
    )
    add_card(
        s,
        Inches(0.45),
        Inches(4.1),
        Inches(6.1),
        Inches(2.4),
        "Retrieval stack used here",
        [
            "Multilingual dense embeddings + BM25 lexical channel",
            "Query rewrite / HyDE; reciprocal rank fusion; cross-encoder re-rank",
            "At most one corrective re-query when top evidence is weak",
        ],
        accent=GREEN_MID,
    )
    add_card(
        s,
        Inches(6.8),
        Inches(4.1),
        Inches(6.05),
        Inches(2.4),
        "Pedagogy: SOLO taxonomy",
        [
            "Biggs & Collis: structural complexity of learner responses",
            "Guides tiered item design from unistructural to relational",
            "Supports practice beyond one-shot QA",
        ],
        accent=GREEN_DARK,
    )
    slides_meta.append((s, "Background Review"))

    # ========== 6 System architecture (tech from PDF) ==========
    s = new_slide(prs)
    add_title_bar(s, "Technical Implementation — Architecture")
    note = s.shapes.add_textbox(Inches(0.45), Inches(1.15), Inches(12.4), Inches(0.4))
    _set_text(
        note.text_frame,
        "Single-process FastAPI + static frontend; embeddings/re-ranking run locally; external LLM used for generation (qwen-plus).",
        size=SZ_LEAD,
        color=MUTED,
    )
    layers = [
        ("Presentation", "Login · Teacher console · Assistant · KB · Quiz · Export", GREEN_LIGHT),
        ("Application", "Auth/admin · Session/file/QA APIs · Exercises · Orchestration", GREEN_MID),
        ("Domain", "Parsing/OCR · Hybrid retrieval · Grounded generation · Quiz/grading", GREEN_DARK),
        ("Infrastructure", "SQLite (auth+KB) · ChromaDB (sessions) · FAISS vectors · Files", GREEN_DEEPER),
    ]
    for i, (name, desc, color) in enumerate(layers):
        top = Inches(1.7 + i * 1.2)
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), top, Inches(12.4), Inches(1.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = CARD
        bar.line.color.rgb = LINE
        tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), top + Inches(0.22), Inches(2.4), Inches(0.6))
        _rgb(tag, color)
        tt = s.shapes.add_textbox(Inches(0.65), top + Inches(0.3), Inches(2.4), Inches(0.5))
        _set_text(tt.text_frame, name, size=SZ_LABEL, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        db = s.shapes.add_textbox(Inches(3.3), top + Inches(0.3), Inches(9.2), Inches(0.55))
        _set_text(db.text_frame, desc, size=SZ_CARD_TITLE, color=TEXT)
    foot = s.shapes.add_textbox(Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.4))
    _set_text(
        foot.text_frame,
        "ChromaDB stores sessions/messages/quizzes — not document vectors. Document vectors live in FAISS.",
        size=SZ_BODY,
        color=MUTED,
    )
    slides_meta.append((s, "Technical Implementation"))

    # ========== 7 Pipeline ==========
    s = new_slide(prs)
    add_title_bar(s, "Technical Implementation — End-to-End Pipeline")
    stages = [
        ("Ingest", ["Parse (+OCR)", "Token chunk", "Parent–child", "Embed → FAISS"]),
        ("Retrieve", ["Rewrite + HyDE", "Dense ∥ BM25", "RRF fusion", "Re-rank ≤1 fix"]),
        ("Generate", ["3-band gate", "Refine / LiM", "LLM answer", "Citation repair"]),
        ("Teach loop", ["SOLO quiz", "Grading", "Teacher export", "Fail-safe fallback"]),
    ]
    for i, (title, steps) in enumerate(stages):
        left = Inches(0.4 + i * 3.2)
        head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), Inches(3.0), Inches(0.55))
        _rgb(head, GREEN_DARK if i % 2 == 0 else GREEN_LIGHT)
        ht = s.shapes.add_textbox(left, Inches(1.38), Inches(3.0), Inches(0.45))
        _set_text(ht.text_frame, title, size=SZ_CARD_TITLE, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, step in enumerate(steps):
            top = Inches(2.1 + j * 1.05)
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.0), Inches(0.85))
            box.fill.solid()
            box.fill.fore_color.rgb = CARD
            box.line.color.rgb = LINE
            tb = s.shapes.add_textbox(left + Inches(0.12), top + Inches(0.22), Inches(2.75), Inches(0.5))
            _set_text(tb.text_frame, step, size=SZ_BODY, color=TEXT, align=PP_ALIGN.CENTER)
        if i < 3:
            ar = s.shapes.add_textbox(left + Inches(2.85), Inches(3.6), Inches(0.4), Inches(0.4))
            _set_text(ar.text_frame, "→", size=18, bold=True, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
    slides_meta.append((s, "Technical Implementation"))

    # ========== 8 Grounding gate detail ==========
    s = new_slide(prs)
    add_title_bar(s, "Three-Band Grounding Gate")
    intro = s.shapes.add_textbox(Inches(0.45), Inches(1.15), Inches(12.4), Inches(0.55))
    intro.text_frame.word_wrap = True
    _set_text(
        intro.text_frame,
        "CRAG-inspired gate. Decision emphasises the second-ranked hit: on this corpus, course-covered questions often have two strong passages; out-of-scope items often show one accidental hit then a drop.",
        size=SZ_LEAD,
        color=MUTED,
    )
    bands = [
        ("none", "Evidence too weak.\n→ labelled general answer\n(not a course conclusion)", RGBColor(0xB9, 0x1C, 0x1C)),
        ("weak", "Partial support.\n→ labelled general answer;\nboundary may call sufficiency judge", RGBColor(0xB4, 0x53, 0x09)),
        ("grounded", "Treated as course-supported.\n→ constrained LLM answer\nwith citations", GREEN_DARK),
    ]
    band_top, band_h, band_w, band_gap = Inches(1.85), Inches(2.55), Inches(4.0), Inches(0.2)
    for i, (name, desc, color) in enumerate(bands):
        left = Inches(0.45) + i * (band_w + band_gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, band_top, band_w, band_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = LINE
        tag = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + Inches(0.25),
            band_top + Inches(0.22),
            Inches(3.5),
            Inches(0.5),
        )
        _rgb(tag, color)
        tt = s.shapes.add_textbox(left + Inches(0.25), band_top + Inches(0.28), Inches(3.5), Inches(0.4))
        _set_text(tt.text_frame, name, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        bd = s.shapes.add_textbox(
            left + Inches(0.25),
            band_top + Inches(0.9),
            Inches(3.5),
            Inches(1.45),
        )
        bd.text_frame.word_wrap = True
        _add_paras(bd.text_frame, desc.split("\n"), size=SZ_BODY, color=TEXT, spacing=4)
    # thresholds — clear gap below bands (bands end ~4.40)
    th = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(4.65), Inches(12.4), Inches(2.25))
    th.fill.solid()
    th.fill.fore_color.rgb = CARD
    th.line.color.rgb = LINE
    tht = s.shapes.add_textbox(Inches(0.7), Inches(4.85), Inches(12.0), Inches(1.9))
    tht.text_frame.word_wrap = True
    _add_paras(
        tht.text_frame,
        [
            "Cosine-regime thresholds calibrated on ELEC6081 (Table 5.3 / §8.5)",
            "multi_strong / single_strong / second_support = 0.62 / 0.75 / 0.62",
            "Re-ranking regime uses a separate threshold family (strong ≈ 0.80); not interchangeable.",
            "Boundary weak cases: optional sufficiency judge on top-5; labelled fallback if judge fails.",
        ],
        size=SZ_LEAD,
        color=TEXT,
        bold_first=True,
        spacing=6,
        sub_size=SZ_BODY,
    )
    slides_meta.append((s, "Technical Implementation"))

    # ========== 9 Benchmark ==========
    s = new_slide(prs)
    add_title_bar(s, "Benchmark Composition")
    add_card(
        s,
        Inches(0.45),
        Inches(1.25),
        Inches(6.1),
        Inches(5.3),
        "ELEC6081 evaluation setting",
        [
            "Course: graduate biomedical signal processing materials",
            "65 questions: on-topic + deliberately out-of-scope",
            "Fields: question, ground_truth, expected_grounding (for OOS)",
            "Harness: tools/rag_eval.py (full and grounding-only modes)",
            "Gate calibration embedding: BAAI/bge-small-zh-v1.5 (re-rank off)",
            "Re-rank regime: multilingual-e5-small + MiniLM cross-encoder",
            "Generation endpoint: DashScope-compatible qwen-plus",
        ],
        accent=GREEN_DARK,
    )
    add_card(
        s,
        Inches(6.8),
        Inches(1.25),
        Inches(6.05),
        Inches(5.3),
        "What is measured",
        [
            "Grounded precision / recall of the three-band gate (65 questions)",
            "Context precision: fusion-order vs re-ranking (10-question probe)",
            "Ownership / export-boundary checks (deterministic unit tests)",
            "Design goal: answer when course-supported; otherwise label as general knowledge",
            "Thresholds are corpus-specific — not transferred from public benchmarks",
            "On-topic and off-topic items are both needed to assess the gate",
        ],
        accent=GREEN_LIGHT,
    )
    slides_meta.append((s, "Benchmark Composition"))

    # ========== 10 Results gate ==========
    s = new_slide(prs)
    add_title_bar(s, "Evaluation and Results — Grounding Gate (H1)")
    # big metrics
    metrics = [
        ("0.982", "Grounded precision", "65-q ELEC6081, cosine"),
        ("0.982", "Grounded recall", "65-q ELEC6081, cosine"),
        ("0.62/0.75/0.62", "Threshold triple", "multi / single / second"),
    ]
    for i, (val, label, sub) in enumerate(metrics):
        left = Inches(0.45 + i * 4.2)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), Inches(4.0), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = LINE
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.3), Inches(4.0), Inches(0.1))
        _rgb(strip, GREEN_LIGHT if i != 2 else GREEN_DARK)
        vb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.6), Inches(3.7), Inches(0.8))
        _set_text(vb.text_frame, val, size=SZ_METRIC, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
        lb = s.shapes.add_textbox(left + Inches(0.15), Inches(2.5), Inches(3.7), Inches(0.7))
        _add_paras(lb.text_frame, [label, sub], size=SZ_CARD_TITLE, color=TEXT, bold_first=True, spacing=2, sub_size=11)
        for j, p in enumerate(lb.text_frame.paragraphs):
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                if j == 1:
                    _apply_font(r, size=11, bold=False, color=MUTED)
    body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(3.8), Inches(12.4), Inches(2.7))
    body.fill.solid()
    body.fill.fore_color.rgb = CARD
    body.line.color.rgb = LINE
    bt = s.shapes.add_textbox(Inches(0.75), Inches(4.05), Inches(11.9), Inches(2.3))
    _add_paras(
        bt.text_frame,
        [
            "Finding (relative to alternatives on this set; not a global optimum claim)",
            "• Raising only the top-hit threshold can cut fabrication but refuses many legitimate course questions.",
            "• Requiring two passages to clear 0.62 yields high precision and recall together here (0.982 / 0.982).",
            "• Selected thresholds are pinned by unit tests against a high-precision-only alternative.",
            "• Fail-safe: weak / none → explicitly labelled general-knowledge answer (not a silent course claim).",
        ],
        size=SZ_LEAD,
        color=TEXT,
        bold_first=True,
        spacing=5,
        sub_size=SZ_BODY,
    )
    slides_meta.append((s, "Evaluation and Results"))

    # ========== 11 Results rerank ==========
    s = new_slide(prs)
    add_title_bar(s, "Evaluation and Results — Re-ranking (H2)")
    # before/after
    for i, (title, val, color) in enumerate(
        [("Fusion-order baseline", "0.380", MUTED), ("After cross-encoder", "0.613", GREEN_DARK)]
    ):
        left = Inches(0.45 + i * 6.3)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), Inches(6.0), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = LINE
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.3), Inches(6.0), Inches(0.1))
        _rgb(strip, GREEN_LIGHT if i == 1 else GREEN_DARK)
        t = s.shapes.add_textbox(left + Inches(0.3), Inches(1.6), Inches(5.4), Inches(0.4))
        _set_text(t.text_frame, title, size=SZ_CARD_TITLE, bold=True, color=MUTED if i == 0 else GREEN_DARK)
        v = s.shapes.add_textbox(left + Inches(0.3), Inches(2.2), Inches(5.4), Inches(0.9))
        _set_text(v.text_frame, f"context_precision = {val}", size=SZ_METRIC, bold=True, color=color)
        s2 = s.shapes.add_textbox(left + Inches(0.3), Inches(3.15), Inches(5.4), Inches(0.35))
        _set_text(s2.text_frame, "mean over 10-question probe", size=11, color=MUTED)

    body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(4.0), Inches(12.4), Inches(2.5))
    body.fill.solid()
    body.fill.fore_color.rgb = CARD
    body.line.color.rgb = LINE
    bt = s.shapes.add_textbox(Inches(0.75), Inches(4.25), Inches(11.9), Inches(2.1))
    _add_paras(
        bt.text_frame,
        [
            "Δ = +0.233 on mean context_precision; higher on 7 of 10 probe questions.",
            "• Example effect: demotes quiz/tutorial fragments that mainly share surface vocabulary.",
            "• Cross-encoder scores query–passage pairs jointly; bi-encoder retrieval scores do not.",
            "• Under the re-rank score scale, gate thresholds need re-calibration (H3 / scale discipline).",
        ],
        size=SZ_LEAD,
        color=TEXT,
        bold_first=True,
        spacing=6,
        sub_size=SZ_BODY,
    )
    slides_meta.append((s, "Evaluation and Results"))

    # ========== 12 Visual fig 4.1 ==========
    s = new_slide(prs)
    add_title_bar(s, "Visual Results — Container View (Fig. 4.1)")
    fig = FIGS / "fig_4_1_architecture.png"
    if fig.exists():
        # keep clear of header/footer: max height ~5.6"
        s.shapes.add_picture(str(fig), Inches(1.55), Inches(1.25), height=Inches(5.6))
    else:
        tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
        _set_text(tb.text_frame, "[Figure 4.1 not found in figures/]", size=16, color=MUTED, align=PP_ALIGN.CENTER)
    slides_meta.append((s, "Visual Results"))

    # ========== 13 Visual fig 5.3 ==========
    s = new_slide(prs)
    add_title_bar(s, "Visual Results — Grounding Gate (Fig. 5.3)")
    fig = FIGS / "fig_5_3_grounding.png"
    if fig.exists():
        s.shapes.add_picture(str(fig), Inches(1.55), Inches(1.25), height=Inches(5.6))
    else:
        tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
        _set_text(tb.text_frame, "[Figure 5.3 not found in figures/]", size=16, color=MUTED, align=PP_ALIGN.CENTER)
    slides_meta.append((s, "Visual Results"))

    # ========== 14 Conclusions ==========
    s = new_slide(prs)
    add_title_bar(s, "Conclusions")
    points = [
        ("Product", "A deployable teacher–student co-learning bot over one shared course KB."),
        ("Method", "Hybrid retrieval + calibrated three-band gate + citation coverage repair."),
        ("Pedagogy", "SOLO-informed practice and an export path that preserves privacy boundaries."),
        ("Evidence", "On ELEC6081: gate P/R ≈ 0.982; re-ranking lifts context precision by +0.233."),
    ]
    for i, (t, b) in enumerate(points):
        top = Inches(1.25 + i * 1.3)
        tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(2.2), Inches(1.0))
        _rgb(tag, GREEN_DARK if i % 2 == 0 else GREEN_LIGHT)
        tt = s.shapes.add_textbox(Inches(0.5), top + Inches(0.28), Inches(2.2), Inches(0.5))
        _set_text(tt.text_frame, t, size=SZ_CARD_TITLE, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.95), top, Inches(9.85), Inches(1.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = LINE
        bb = s.shapes.add_textbox(Inches(3.2), top + Inches(0.28), Inches(9.4), Inches(0.55))
        _set_text(bb.text_frame, b, size=SZ_CARD_TITLE, color=TEXT)
    slides_meta.append((s, "Conclusions"))

    # ========== 15 Limitations & future ==========
    s = new_slide(prs)
    add_title_bar(s, "Limitations and Future Work")
    add_card(
        s,
        Inches(0.45),
        Inches(1.25),
        Inches(6.1),
        Inches(5.3),
        "Limitations",
        [
            "Single course / single instance — no multi-course isolation yet",
            "Gate thresholds do not transfer across models or growing KBs",
            "Complex layouts, noisy scans, and estimated page numbers remain weak",
            "Practice adaptivity is still shallow compared with full tutoring",
        ],
        accent=GREEN_DARK,
    )
    add_card(
        s,
        Inches(6.8),
        Inches(1.25),
        Inches(6.05),
        Inches(5.3),
        "Future work",
        [
            "Multi-course configuration and stronger campus deployment (HTTPS / IdP)",
            "Better layout-aware parsing and more reliable location anchors",
            "Continued recalibration of gate and re-rank thresholds on course sets",
            "Deeper SOLO-adaptive practice and richer teacher analytics",
        ],
        accent=GREEN_LIGHT,
    )
    slides_meta.append((s, "Limitations and Future Work"))

    # ========== 16 Thank you ==========
    s = new_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    _rgb(bg, GREEN_DARK)
    add_logo(s)
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.2))
    _set_text(t.text_frame, "Thank you", size=SZ_THANK, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.6))
    _set_text(t2.text_frame, "Questions & Discussion", size=20, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
    t3 = s.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.0))
    _add_paras(
        t3.text_frame,
        ["Jin Yue  ·  3036494489", "Supervisor: Dr. Andrew H.C. Wu  ·  ECE, HKU"],
        size=SZ_COVER_SUB,
        color=WHITE,
        spacing=6,
    )
    for p in t3.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    slides_meta.append((s, ""))

    # footers (skip title & thank-you full-bleed)
    total = len(slides_meta)
    for i, (slide, section) in enumerate(slides_meta, 1):
        if i in (1, total):
            continue
        add_footer(slide, i, total, section)

    # save as NEW file only
    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()
    dest = OUT_DIR / "Jin_Yue_3036494489_Final_Defense_Clean_TNR.pptx"
    try:
        dest.write_bytes(data)
    except PermissionError:
        dest = OUT_DIR / f"Jin_Yue_3036494489_Final_Defense_Clean_TNR_{int(time.time())}.pptx"
        dest.write_bytes(data)
    # also try updating Clean.pptx
    clean = OUT_DIR / OUT_NAME
    try:
        clean.write_bytes(data)
    except PermissionError:
        pass
    bak = Path(r"e:\For_RAG\docs") / dest.name
    try:
        bak.write_bytes(data)
    except Exception:
        bak = None
    print(f"Wrote {dest} ({dest.stat().st_size} bytes), slides={total}")
    if bak:
        print(f"Backup {bak}")
    return dest


if __name__ == "__main__":
