#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build final dissertation defense PPT from HKU.pptx template,
aligned with the requested framework and thesis/mid-term content.
"""
from __future__ import annotations

import copy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

TEMPLATE = Path(r"C:\Users\10603\Desktop\毕业论文模版\HKU.pptx")
OUT_DIR = Path(r"C:\Users\10603\Desktop\毕业论文模版")
OUT_WORK = Path(r"e:\For_RAG\tools\_final_defense_build.pptx")
OUT = OUT_DIR / "Jin_Yue_3036494489_Final_Defense.pptx"
FIGS = Path(r"C:\Users\10603\Desktop\毕业论文模版\第十一稿\figures")

HKU_GREEN = RGBColor(0x27, 0x58, 0x29)
DARK = RGBColor(0x1F, 0x2A, 0x2E)
GRAY = RGBColor(0x4A, 0x55, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_runs_text(shape, text: str, *, size=None, bold=None, color=None, clear_extra=True):
    """Replace text in a shape, keeping first paragraph/run style when possible."""
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    lines = text.split("\n") if text is not None else [""]
    # Ensure enough paragraphs
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for i, para in enumerate(tf.paragraphs):
        if i >= len(lines):
            if clear_extra:
                for r in para.runs:
                    r.text = ""
            continue
        line = lines[i]
        if not para.runs:
            run = para.add_run()
            run.text = line
            if size is not None:
                run.font.size = size
            if bold is not None:
                run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color
        else:
            para.runs[0].text = line
            if size is not None:
                para.runs[0].font.size = size
            if bold is not None:
                para.runs[0].font.bold = bold
            if color is not None:
                try:
                    para.runs[0].font.color.rgb = color
                except Exception:
                    pass
            for r in para.runs[1:]:
                r.text = ""


def text_shapes(slide):
    out = []
    for sh in slide.shapes:
        if getattr(sh, "has_text_frame", False):
            t = "\n".join(p.text for p in sh.text_frame.paragraphs).strip()
            out.append((sh, t))
    return out


def delete_slide(prs: Presentation, index: int) -> None:
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def clear_placeholder_bullets(shape, bullets: list[str], *, size=Pt(14)):
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        for r in p.runs:
            r.font.size = size
            r.font.color.rgb = DARK
            r.font.name = "Calibri"


def fill_by_order(slide, replacements: list[str], *, skip_nav=True):
    """
    Fill text shapes top-to-bottom / left-to-right, skipping bottom nav labels
    when skip_nav and text matches known nav words.
    """
    nav = {
        "introduction",
        "literature review",
        "methods",
        "results",
        "conclusion",
        "overview",
    }
    shapes = []
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        t = "\n".join(p.text for p in sh.text_frame.paragraphs).strip().lower()
        if skip_nav and t in nav:
            continue
        # skip tiny numbering boxes that are only digits
        if t.isdigit() and len(t) <= 2:
            continue
        shapes.append(sh)
    shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    for sh, txt in zip(shapes, replacements):
        set_runs_text(sh, txt)


def update_bottom_nav(slide, labels: list[str]):
    """Update the five bottom navigation labels if present."""
    nav_order = [
        "introduction",
        "literature review",
        "methods",
        "results",
        "conclusion",
    ]
    found = {k: None for k in nav_order}
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        t = "\n".join(p.text for p in sh.text_frame.paragraphs).strip().lower()
        if t in found:
            found[t] = sh
    # map in template order
    mapping = list(zip(nav_order, labels))
    for old, new in mapping:
        sh = found.get(old)
        if sh is not None and new:
            set_runs_text(sh, new, size=Pt(10), bold=True)


NAV = ["Scope", "Background", "Benchmark", "Results", "Conclusion"]


def main():
    # Load template directly (do not overwrite a locked work file)
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)

    # ---------- Slide 1: Title ----------
    s = slides[0]
    for sh, t in text_shapes(s):
        low = t.lower()
        if "presentation title" in low or "goes here" in low:
            set_runs_text(
                sh,
                "Design and Implementation of a\nRAG-Based Teacher–Student\nCo-Learning System",
                size=Pt(28),
                bold=True,
                color=HKU_GREEN,
            )
        elif "subtopic" in low or "school" in low:
            set_runs_text(
                sh,
                "Department of Electrical and Electronic Engineering\nThe University of Hong Kong",
                size=Pt(14),
                color=DARK,
            )
        elif "dissertation defense" in low or "your name" in low or "month" in low:
            set_runs_text(
                sh,
                "Final Defense  |  July 2026  |  Jin Yue (3036494489)",
                size=Pt(14),
                bold=True,
            )

    # ---------- Slide 2: Overview ----------
    s = slides[1]
    # numbered boxes + labels
    for sh, t in text_shapes(s):
        if t.strip().upper() == "OVERVIEW":
            set_runs_text(sh, "OVERVIEW", size=Pt(28), bold=True, color=HKU_GREEN)
        elif t.strip() == "1":
            set_runs_text(sh, "1")
        elif t.strip() == "2":
            set_runs_text(sh, "2")
        elif t.strip() == "3":
            set_runs_text(sh, "3")
        elif t.strip() == "4":
            set_runs_text(sh, "4")
        elif t.strip() == "5":
            set_runs_text(sh, "5")
        elif t.strip() == "Introduction":
            set_runs_text(sh, "Task Scope &\nMotivation")
        elif t.strip() == "Literature Review":
            set_runs_text(sh, "Background\nReview")
        elif t.strip() == "Research Methods":
            set_runs_text(sh, "Benchmark &\nEvaluation")
        elif t.strip() == "Results":
            set_runs_text(sh, "Visual\nResults")
        elif t.strip() == "Conclusion":
            set_runs_text(sh, "Conclusions &\nFuture Work")

    # ---------- Slide 3: Task Scope ----------
    s = slides[2]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        low = t.lower()
        if "background" in low and "briefly" in low:
            set_runs_text(
                sh,
                "Task Scope: build a course-deployable Teacher–Student Co-Learning (SOLO) Bot",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
        elif t.strip() == "Add your title" or "add your title" in low:
            # will set by order below
            pass
    # Fill the three cards by position
    cards_title = []
    cards_body = []
    for sh, t in text_shapes(s):
        if t.strip() == "Add your title":
            cards_title.append(sh)
        elif "text can be changed" in t.lower():
            cards_body.append(sh)
    titles = [
        "Shared Course KB",
        "Traceable RAG QA",
        "SOLO Practice Loop",
    ]
    bodies = [
        "One course knowledge base written by teachers and read by students; supports PDF/PPT/Word/Excel/images with OCR.",
        "Hybrid retrieval + grounding gate + sentence-level citations linking answers to file and page/slide.",
        "Tiered exercise generation/grading under SOLO, plus teacher export of questions and scores (no full transcripts).",
    ]
    for sh, txt in zip(cards_title, titles):
        set_runs_text(sh, txt, size=Pt(14), bold=True, color=WHITE)
    for sh, txt in zip(cards_body, bodies):
        set_runs_text(sh, txt, size=Pt(12), color=DARK)

    # ---------- Slide 4: Motivation ----------
    s = slides[3]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        low = t.lower()
        if "significance" in low or "important" in low:
            set_runs_text(
                sh,
                "Motivation: why a course RAG system is needed",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
        elif "theoretical significance" in low:
            set_runs_text(sh, "Core Problems", size=Pt(14), bold=True, color=HKU_GREEN)
        elif "practical significance" in low:
            set_runs_text(sh, "Design Response", size=Pt(14), bold=True, color=HKU_GREEN)
        elif "text can be changed" in low:
            pass
    bodies = []
    for sh, t in text_shapes(s):
        if "text can be changed" in t.lower():
            bodies.append(sh)
    if len(bodies) >= 2:
        clear_placeholder_bullets(
            bodies[0],
            [
                "General LLMs cannot sync with private course materials",
                "Answers often lack locatable evidence (hallucination risk)",
                "QA alone does not close the teaching–learning loop",
                "Course corpora are multi-format and bilingual",
            ],
            size=Pt(13),
        )
        clear_placeholder_bullets(
            bodies[1],
            [
                "Shared KB + hybrid retrieval (dense + BM25 + RRF)",
                "Three-band grounding gate + citation coverage repair",
                "SOLO-informed practice and feedback export",
                "Single-process FastAPI deployment for one course",
            ],
            size=Pt(13),
        )

    # ---------- Slide 5: Background Review ----------
    s = slides[4]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "summarize relevant" in t.lower() or "literature" in t.lower() and "summarize" in t.lower():
            set_runs_text(
                sh,
                "Background Review: RAG for trustworthy course QA",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
        elif "text can be changed" in t.lower():
            pass
    bodies = [sh for sh, t in text_shapes(s) if "text can be changed" in t.lower()]
    texts = [
        "RAG combines parametric LLMs with non-parametric retrieval so answers can be grounded in updatable corpora (Lewis et al.).",
        "Education settings need private materials, verifiable citations, and more than one-shot QA — a closed instructional loop.",
        "Trustworthiness requires gates, refinement, and citation checks; retrieving more passages alone is insufficient.",
    ]
    for sh, txt in zip(bodies, texts):
        set_runs_text(sh, txt, size=Pt(13), color=DARK)

    # ---------- Slide 6: Background highlights ----------
    s = slides[5]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "hightlights" in t.lower() or "highlights" in t.lower():
            set_runs_text(
                sh,
                "Background Review: methods used in this system",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
    titles = [sh for sh, t in text_shapes(s) if t.strip() == "Add title"]
    bodies = [sh for sh, t in text_shapes(s) if "text can be changed" in t.lower()]
    t4 = ["Hybrid Retrieval", "Grounding Gate", "SOLO Taxonomy", "Parent–Child Chunks"]
    b4 = [
        "Dense + BM25 fused by RRF, then cross-encoder re-ranking with graceful fallback.",
        "Three bands: none / weak / grounded; decide mainly on the second hit.",
        "Tiered practice from Biggs & Collis; aligns generation with cognitive levels.",
        "Retrieve small children; expand to page/slide parents for generation.",
    ]
    for sh, txt in zip(titles, t4):
        set_runs_text(sh, txt, size=Pt(13), bold=True, color=HKU_GREEN)
    for sh, txt in zip(bodies, b4):
        set_runs_text(sh, txt, size=Pt(11), color=DARK)

    # ---------- Slide 7: Gaps → Task framing ----------
    s = slides[6]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "gaps" in t.lower():
            set_runs_text(
                sh,
                "From gaps to system requirements",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
    titles = [sh for sh, t in text_shapes(s) if t.strip() == "Add title"]
    bodies = [sh for sh, t in text_shapes(s) if "text can be changed" in t.lower()]
    t4 = ["Private corpus", "Traceability", "Learning loop", "Deployability"]
    b4 = [
        "Must index teacher materials and session uploads with OCR fallback.",
        "Every course claim must map to file + page/slide; weak evidence → labelled fallback.",
        "QA → practice → grading → teacher feedback export.",
        "Single-worker FastAPI + Docker; local embeddings/re-rank; hosted generation only.",
    ]
    for sh, txt in zip(titles, t4):
        set_runs_text(sh, txt, size=Pt(13), bold=True, color=HKU_GREEN)
    for sh, txt in zip(bodies, b4):
        set_runs_text(sh, txt, size=Pt(11), color=DARK)

    # ---------- Slide 8: Pipeline / theoretical basis layout ----------
    s = slides[7]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        low = t.lower()
        if "theoretical basis" in low:
            set_runs_text(
                sh,
                "System pipeline (end-to-end)",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
        elif t.strip() == "Add title":
            set_runs_text(
                sh,
                "Ingest → Retrieve → Gate → Generate → Cite → Practice/Export",
                size=Pt(14),
                bold=True,
            )
        elif "text can be changed" in low:
            set_runs_text(
                sh,
                "Parse (+OCR) → token-budget chunks → embed/FAISS;\n"
                "Rewrite+HyDE → Dense∥BM25 → RRF → re-rank → corrective ≤1;\n"
                "Three-band gate → refine/LiM order → LLM answer → citation repair;\n"
                "SOLO quizzes & teacher export (aggregated only).",
                size=Pt(13),
            )
        elif t.strip().upper() == "START":
            set_runs_text(sh, "UPLOAD\n/ ASK")
        elif t.strip().upper() == "END":
            set_runs_text(sh, "ANSWER\n+ QUIZ")
        elif t.strip() == "1":
            set_runs_text(sh, "1")
        elif t.strip() == "2":
            set_runs_text(sh, "2")

    # ---------- Slides 9–11: Evaluation results (reuse results layouts) ----------
    # Slide 9
    s = slides[8]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "present your key findings" in t.lower():
            set_runs_text(
                sh,
                "Benchmark Composition (ELEC6081)",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
        elif t.strip() == "Chart title" or t.strip() == "Add title":
            set_runs_text(sh, "Evaluation set design", size=Pt(14), bold=True)
        elif "text can be changed" in t.lower():
            set_runs_text(
                sh,
                "• 65 course questions (on-topic + out-of-scope)\n"
                "• Corpus: course materials (calibration and regression)\n"
                "• Metrics: grounded precision/recall; context precision\n"
                "• Tools: tools/rag_eval.py (full & grounding-only)\n"
                "• Models: bge-small-zh (gate calib.); e5-small + MiniLM (re-rank regime)",
                size=Pt(13),
            )

    # Slide 10 — two column findings
    s = slides[9]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "present your key findings" in t.lower():
            set_runs_text(
                sh,
                "Evaluation and Results — Grounding Gate (H1)",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
    titles = [sh for sh, t in text_shapes(s) if t.strip() == "Add title"]
    bodies = [sh for sh, t in text_shapes(s) if "text can be changed" in t.lower()]
    if titles:
        set_runs_text(titles[0], "Selected thresholds (cosine)", size=Pt(13), bold=True)
    if len(titles) > 1:
        set_runs_text(titles[1], "Key finding", size=Pt(13), bold=True)
    if bodies:
        set_runs_text(
            bodies[0],
            "multi_strong / single_strong / second_support\n"
            "= 0.62 / 0.75 / 0.62\n\n"
            "Grounded precision ≈ 0.982\n"
            "Grounded recall ≈ 0.982",
            size=Pt(13),
        )
    if len(bodies) > 1:
        set_runs_text(
            bodies[1],
            "Thresholding the second hit outperforms merely raising the top-hit bar:\n"
            "high precision without collapsing recall.\n"
            "Fail-safe: weak/none → labelled general answer.",
            size=Pt(13),
        )

    # Slide 11
    s = slides[10]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "present your key findings" in t.lower():
            set_runs_text(
                sh,
                "Evaluation and Results — Re-ranking (H2)",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
        elif t.strip() == "Add title":
            set_runs_text(sh, "Context precision on 10-question probe", size=Pt(14), bold=True)
        elif "text can be changed" in t.lower():
            set_runs_text(
                sh,
                "Fusion-order baseline → after cross-encoder re-ranking:\n"
                "mean context_precision 0.380 → 0.613 (Δ +0.233).\n"
                "Rises on 7/10 questions; demotes quiz fragments that only share vocabulary.\n"
                "Under re-rank score scale, thresholds must be re-calibrated (H3).",
                size=Pt(13),
            )

    # ---------- Slide 12: Visual Results — architecture ----------
    s = slides[11]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "present your key findings" in t.lower():
            set_runs_text(
                sh,
                "Visual Results — System container view",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
        elif t.strip() == "Add title":
            set_runs_text(sh, "Figure 4.1 Deployed system", size=Pt(12), bold=True)
        elif "text can be changed" in t.lower():
            set_runs_text(
                sh,
                "Browser → FastAPI single process → local stores (SQLite/Chroma/FAISS/files);\n"
                "only generation may call external LLM (qwen-plus).",
                size=Pt(12),
            )
    # add figure if space: place on rightish area
    fig = FIGS / "fig_4_1_architecture.png"
    if fig.exists():
        # avoid overlapping title: lower-middle
        try:
            s.shapes.add_picture(str(fig), Inches(1.2), Inches(2.0), width=Inches(10.5))
        except Exception:
            pass

    # ---------- Slide 13: Visual — pipeline / gate ----------
    s = slides[12]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "present your key findings" in t.lower():
            set_runs_text(
                sh,
                "Visual Results — Grounding gate & pipeline",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
        elif t.strip() == "Add title":
            set_runs_text(sh, "Figures 5.1 / 5.3", size=Pt(12), bold=True)
        elif "text can be changed" in t.lower():
            set_runs_text(
                sh,
                "End-to-end ingest–retrieve–generate pipeline;\n"
                "three-band gate: none / weak / grounded.",
                size=Pt(12),
            )
    fig = FIGS / "fig_5_3_grounding.png"
    if fig.exists():
        try:
            s.shapes.add_picture(str(fig), Inches(1.2), Inches(2.0), width=Inches(10.5))
        except Exception:
            pass

    # ---------- Slide 14: Summary findings ----------
    s = slides[13]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "summary of key findings" in t.lower():
            set_runs_text(
                sh,
                "Conclusions — what was delivered",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
        elif t.strip() == "Add your title":
            set_runs_text(sh, "Three contributions", size=Pt(14), bold=True, color=WHITE)
    bodies = [sh for sh, t in text_shapes(s) if "text can be changed" in t.lower()]
    if bodies:
        clear_placeholder_bullets(
            bodies[0],
            [
                "Shared course KB with multi-format parsing and OCR",
                "Traceable RAG: hybrid retrieval, grounding gate, citations",
                "SOLO practice + teacher feedback export boundary",
                "Validated on ELEC6081 questions; gated P/R ≈ 0.982",
            ],
            size=Pt(13),
        )
    if len(bodies) > 1:
        clear_placeholder_bullets(
            bodies[1],
            [
                "Deployable prototype: FastAPI + Docker + local models",
                "Fail-safe behaviour when evidence is weak",
                "Deterministic tests for ownership and export rules",
            ],
            size=Pt(13),
        )

    # ---------- Slide 15: dual summary ----------
    s = slides[14]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "summary of key findings" in t.lower():
            set_runs_text(
                sh,
                "Conclusions — experimental takeaways",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
    titles = [sh for sh, t in text_shapes(s) if t.strip() == "Add your title"]
    bodies = [sh for sh, t in text_shapes(s) if "text can be changed" in t.lower()]
    if titles:
        set_runs_text(titles[0], "Gate calibration", size=Pt(13), bold=True)
    if len(titles) > 1:
        set_runs_text(titles[1], "Re-ranking", size=Pt(13), bold=True)
    if bodies:
        set_runs_text(
            bodies[0],
            "Second-hit thresholding yields high precision and recall together; raising only the top hit trades away too much recall.",
            size=Pt(12),
        )
    if len(bodies) > 1:
        set_runs_text(
            bodies[1],
            "Cross-encoder improves context precision (+0.233 on the probe); score-scale change requires re-calibration.",
            size=Pt(12),
        )

    # ---------- Slide 16: three takeaways ----------
    s = slides[15]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "summary of key findings" in t.lower():
            set_runs_text(
                sh,
                "Conclusions — operational principles",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
    titles = [sh for sh, t in text_shapes(s) if t.strip() == "Add title"]
    bodies = [sh for sh, t in text_shapes(s) if "text can be changed" in t.lower()]
    t3 = ["Evidence first", "Scale discipline", "Closed loop"]
    b3 = [
        "Prefer two supporting passages over one lucky top hit.",
        "Cosine and re-rank scores use separate threshold families.",
        "QA, practice, and teacher export form one instructional cycle.",
    ]
    for sh, txt in zip(titles, t3):
        set_runs_text(sh, txt, size=Pt(13), bold=True, color=HKU_GREEN)
    for sh, txt in zip(bodies, b3):
        set_runs_text(sh, txt, size=Pt(12), color=DARK)

    # ---------- Slide 17: Limitations ----------
    s = slides[16]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "limitations" in t.lower():
            set_runs_text(
                sh,
                "Limitations and Future Work",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
    titles = [sh for sh, t in text_shapes(s) if t.strip() == "Add your title"]
    bodies = [sh for sh, t in text_shapes(s) if "text can be changed" in t.lower()]
    t4 = ["Single course", "Threshold transfer", "Parsing limits", "Future work"]
    b4 = [
        "One course / one instance; multi-course isolation out of scope.",
        "Gates depend on corpus and embedding/re-ranker; must re-calibrate.",
        "Complex tables/layouts and estimated page numbers remain weak.",
        "Stronger parsing, multi-course config, deeper SOLO adaptivity, HTTPS/IdP.",
    ]
    for sh, txt in zip(titles, t4):
        set_runs_text(sh, txt, size=Pt(12), bold=True, color=WHITE)
    for sh, txt in zip(bodies, b4):
        set_runs_text(sh, txt, size=Pt(11), color=DARK)

    # ---------- Slide 18: Contributions ----------
    s = slides[17]
    update_bottom_nav(s, NAV)
    for sh, t in text_shapes(s):
        if "contribution" in t.lower():
            set_runs_text(
                sh,
                "Contributions to practice",
                size=Pt(16),
                bold=True,
                color=HKU_GREEN,
            )
    titles = [sh for sh, t in text_shapes(s) if t.strip() == "Add your title"]
    bodies = [sh for sh, t in text_shapes(s) if "text can be changed" in t.lower()]
    t4 = ["Product", "Method", "Evaluation", "Engineering"]
    b4 = [
        "Teacher–student co-learning bot over a shared course KB.",
        "Calibrated three-band gate + hybrid retrieval + citations.",
        "Course-specific benchmark with on/off-topic questions.",
        "Tested ownership, export boundary, and deployable stack.",
    ]
    for sh, txt in zip(titles, t4):
        set_runs_text(sh, txt, size=Pt(12), bold=True, color=HKU_GREEN)
    for sh, txt in zip(bodies, b4):
        set_runs_text(sh, txt, size=Pt(11), color=DARK)

    # ---------- Slide 19: Gratitude ----------
    s = slides[18]
    for sh, t in text_shapes(s):
        low = t.lower()
        if "gratitude" in low:
            set_runs_text(sh, "Acknowledgements", size=Pt(24), bold=True, color=HKU_GREEN)
        elif "kim roberts" in low and "ed.d" in low:
            set_runs_text(sh, "Dr. Andrew H.C. Wu", size=Pt(16), bold=True)
        elif "kim roberts" in low:
            set_runs_text(sh, "Department of ECE, HKU", size=Pt(14))
        elif "committee chair" in low:
            set_runs_text(sh, "Supervisor", size=Pt(12))

    # ---------- Slide 20: Thank you ----------
    s = slides[19]
    for sh, t in text_shapes(s):
        low = t.lower()
        if "thank you" in low:
            set_runs_text(sh, "Thank you\nQuestions & Discussion", size=Pt(36), bold=True, color=HKU_GREEN)
        elif "dissertation defense" in low or "your name" in low or "month" in low:
            set_runs_text(
                sh,
                "Final Defense  |  July 2026  |  Jin Yue (3036494489)",
                size=Pt(14),
                bold=True,
            )

    # ---------- Delete color-guide slide (last) ----------
    while len(list(prs.slides)) > 20:
        delete_slide(prs, len(list(prs.slides)) - 1)

    # Save via unique temp path (avoid locked filenames)
    import io
    import time

    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()
    dest = OUT_DIR / f"Jin_Yue_3036494489_Final_Defense.pptx"
    try:
        dest.write_bytes(data)
    except PermissionError:
        dest = OUT_DIR / f"Jin_Yue_3036494489_Final_Defense_{int(time.time())}.pptx"
        dest.write_bytes(data)
    work = Path(r"e:\For_RAG\docs") / dest.name
    try:
        work.write_bytes(data)
    except Exception:
        work = None
    print(f"Wrote {dest} ({dest.stat().st_size} bytes), slides={len(list(prs.slides))}")
    if work:
        print(f"Also: {work}")


if __name__ == "__main__":
    main()
